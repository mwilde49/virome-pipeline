"""
generate_presentation_2026-04-08.py
Lab presentation: DRG Virome Pipeline — Why, How, and What We Found
2026-04-08 — Matthew Wild, TJP Lab, UT Dallas

Updated for v1.5.0: Parkinson's DRG cohort, HSV-1 Tier 1 detection,
BLAST verification offshoot, and future directions (lab-wide data consolidation).

Run from repo root:
    python3 research/presentations/generate_presentation_2026-04-08.py

Outputs:
    research/presentations/virome_lab_presentation_2026-04-08.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
REPO_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
FIG_DIR = os.path.join(REPO_ROOT, "research", "paper1", "figures")
OUT_PATH = os.path.join(
    REPO_ROOT, "research", "presentations",
    "virome_lab_presentation_2026-04-08.pptx",
)

FIG1 = os.path.join(FIG_DIR, "fig1_pipeline_diagram.png")
FIG2 = os.path.join(FIG_DIR, "fig2_filtering_funnel.png")
FIG3 = os.path.join(FIG_DIR, "fig3_tier_summary.png")
FIG4 = os.path.join(FIG_DIR, "fig4_herv_k.png")

# ---------------------------------------------------------------------------
# Design constants (16:9 = 13.33" x 7.5")
# ---------------------------------------------------------------------------
W = Inches(13.33)
H = Inches(7.5)

# Color palette
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
C_BLUE     = RGBColor(0x2E, 0x86, 0xAB)
C_CORAL    = RGBColor(0xE8, 0x48, 0x55)
C_GREEN    = RGBColor(0x3B, 0xB2, 0x73)
C_AMBER    = RGBColor(0xF9, 0xC7, 0x4F)
C_PURPLE   = RGBColor(0x7B, 0x2D, 0x8B)
C_LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)
C_MGRAY    = RGBColor(0x9E, 0x9E, 0x9E)
C_STRIPE   = RGBColor(0xE8, 0xF4, 0xFD)

TITLE_FONT = "Calibri"
BODY_FONT  = "Calibri"
CODE_FONT  = "Courier New"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def rect(slide, x, y, w, h, fill=None, line=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line
        shape.line.width = line_width
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_name=BODY_FONT, font_size=Pt(14), bold=False,
                 italic=False, color=C_DARK, align=PP_ALIGN.LEFT,
                 word_wrap=True, line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if line_spacing:
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txBox


def add_multiline(slide, lines, x, y, w, h,
                  font_name=BODY_FONT, font_size=Pt(13), bold=False,
                  color=C_DARK, line_gap_pt=6, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, b, col = item, bold, color
        else:
            txt = item[0]
            b = item[1] if len(item) > 1 else bold
            col = item[2] if len(item) > 2 else color
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        pPr = p._p.get_or_add_pPr()
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', str(int(line_gap_pt * 100)))
        run = p.add_run()
        run.text = txt
        run.font.name = font_name
        run.font.size = font_size
        run.font.bold = b
        run.font.color.rgb = col
    return txBox


def slide_header(slide, title, subtitle=None,
                 bar_color=C_BLUE, title_color=C_DARK):
    rect(slide, 0, 0, W, Inches(0.08), fill=bar_color)
    add_text_box(slide, title,
                 x=Inches(0.45), y=Inches(0.12), w=Inches(12.0), h=Inches(0.7),
                 font_size=Pt(28), bold=True, color=title_color)
    if subtitle:
        add_text_box(slide, subtitle,
                     x=Inches(0.45), y=Inches(0.78), w=Inches(12.0), h=Inches(0.35),
                     font_size=Pt(14), color=C_MGRAY, italic=True)
    rect(slide, 0, Inches(1.08), W, Inches(0.025), fill=C_LGRAY)


def section_divider(prs, number, title, subtitle, color=C_BLUE):
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, fill=color)
    pill = slide.shapes.add_shape(1, Inches(0.5), Inches(2.8), Inches(1.1), Inches(0.55))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    pill.line.fill.background()
    add_text_box(slide, f"0{number}", Inches(0.5), Inches(2.78), Inches(1.1), Inches(0.6),
                 font_size=Pt(22), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, title,
                 Inches(0.5), Inches(3.45), Inches(12.0), Inches(1.2),
                 font_size=Pt(46), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide, subtitle,
                 Inches(0.5), Inches(4.65), Inches(12.0), Inches(0.7),
                 font_size=Pt(18), color=RGBColor(0xFF, 0xFF, 0xFF), italic=True)
    return slide


def add_figure(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        r = rect(slide, x, y, w, h or Inches(3), fill=C_LGRAY, line=C_MGRAY)
        add_text_box(slide, f"[Figure: {os.path.basename(path)}]",
                     x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), Inches(0.5),
                     font_size=Pt(11), color=C_MGRAY)
        return
    if h:
        slide.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        slide.shapes.add_picture(path, x, y, width=w)


def bullet_box(slide, items, x, y, w, h, font_size=Pt(13.5), gap=8,
               icon="\u25b8", color=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            txt, level = item, 0
        else:
            txt, level = item[0], item[1]
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', str(gap * 100))
        indent_str = "  " * level
        run = p.add_run()
        prefix = "  \u25e6 " if level > 0 else f"{icon} "
        run.text = indent_str + prefix + txt
        run.font.name = BODY_FONT
        run.font.size = font_size
        run.font.color.rgb = (
            color if color is not None
            else (C_DARK if level == 0 else C_MGRAY)
        )
    return txBox


def stat_box(slide, value, label, x, y, w=Inches(2.0), h=Inches(1.2),
             val_color=C_BLUE, bg=C_LGRAY):
    rect(slide, x, y, w, h, fill=bg)
    add_text_box(slide, value, x, y + Inches(0.1), w, Inches(0.65),
                 font_size=Pt(30), bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, x, y + Inches(0.7), w, Inches(0.45),
                 font_size=Pt(11), color=C_DARK, align=PP_ALIGN.CENTER, word_wrap=True)


def footnote(slide, text, y=Inches(7.1)):
    add_text_box(slide, text, Inches(0.45), y, Inches(12.4), Inches(0.35),
                 font_size=Pt(9), color=C_MGRAY, italic=True)


# ---------------------------------------------------------------------------
# Build slides
# ---------------------------------------------------------------------------

def build(prs):

    # ======================================================================
    # SLIDE 1: TITLE
    # ======================================================================
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, Inches(5.2), fill=C_DARK)
    rect(slide, 0, Inches(5.2), W, Inches(2.3), fill=C_BLUE)

    add_text_box(slide,
        "Profiling the Human DRG Virome\nfrom Bulk RNA-seq",
        Inches(0.6), Inches(0.7), Inches(10.5), Inches(2.4),
        font_size=Pt(42), bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide,
        "Why it exists, how it works, and what we've found so far",
        Inches(0.6), Inches(3.1), Inches(10.5), Inches(0.6),
        font_size=Pt(18), color=C_AMBER)
    add_text_box(slide,
        "Matthew Wild  \u00b7  TJP Lab  \u00b7  UT Dallas  \u00b7  April 2026",
        Inches(0.6), Inches(3.75), Inches(10.5), Inches(0.5),
        font_size=Pt(14), color=RGBColor(0xCC, 0xCC, 0xCC))

    # Version badge
    add_text_box(slide,
        "virome-pipeline v1.5.0",
        Inches(10.3), Inches(0.8), Inches(2.8), Inches(0.4),
        font_size=Pt(12), color=C_AMBER, bold=True, align=PP_ALIGN.RIGHT)

    add_text_box(slide,
        "Outline:",
        Inches(0.6), Inches(5.35), Inches(5.5), Inches(0.4),
        font_size=Pt(13), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    items = [
        "Why — neurotropic viruses, the DRG, and an uncharacterized virome",
        "How — dual-database competitive classification from existing RNA-seq",
        "Results — two cohorts, 36 donors, and the first Tier 1 detection",
        "Future — compiling the lab's data, beginning with bulk",
    ]
    for i, item in enumerate(items):
        add_text_box(slide,
            f"{chr(0x2460 + i)}  {item}",
            Inches(0.6), Inches(5.75 + i * 0.38), Inches(11.5), Inches(0.38),
            font_size=Pt(13), color=RGBColor(0xFF, 0xFF, 0xFF))


    # ======================================================================
    # SECTION 1: WHY
    # ======================================================================
    section_divider(prs, 1,
        "Why",
        "The DRG is a natural viral reservoir \u2014 and nobody has systematically profiled it",
        color=C_PURPLE)


    # ── SLIDE 2: WHY VIRUSES IN DRG ─────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Why look for viruses in the dorsal root ganglion?",
                 subtitle="Neurotropic viruses establish latency in sensory ganglia \u2014 and may drive pain long after acute infection",
                 bar_color=C_PURPLE)

    col_w = Inches(3.9)
    col_h = Inches(4.8)
    cols = [
        {
            "title": "The Biology",
            "color": C_BLUE,
            "items": [
                "DRG = primary sensory neurons; gateway for all somatosensory input",
                "VZV establishes latency in DRG \u2014 reactivation \u2192 shingles / PHN",
                "HSV-1/2 latent in trigeminal & sacral DRG",
                "CMV, HHV-6 detected in neural tissue under immunosuppression",
                "DRG has interferon-response satellite glial cells \u2014 constitutive immune surveillance",
            ]
        },
        {
            "title": "The Clinical Need",
            "color": C_CORAL,
            "items": [
                "No systematic profiling of the DRG virome has been reported",
                "Viral triggers implicated in PHN, diabetic neuropathy, HIV neuropathy",
                "HSV-1 linked to neurodegeneration (Parkinson's, Alzheimer's)",
                "HERV-K transcriptional activation in ALS motor neurons",
                "Antiviral therapy associated with reduced PD risk in registry studies",
            ]
        },
        {
            "title": "The Opportunity",
            "color": C_GREEN,
            "items": [
                "Bulk RNA-seq from DRG already exists in this lab \u2014 100\u2013200 donors",
                "Computational pipeline can screen all of it with one command",
                "No new sequencing required for the first virome atlas pass",
                "Establish baseline \u2192 compare disease cohorts \u2192 identify viral associations",
            ]
        }
    ]
    for i, col in enumerate(cols):
        x = Inches(0.3 + i * 4.5)
        rect(slide, x, Inches(1.25), col_w, Inches(0.45), fill=col["color"])
        add_text_box(slide, col["title"],
                     x + Inches(0.1), Inches(1.27), col_w - Inches(0.2), Inches(0.4),
                     font_size=Pt(14), bold=True,
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        rect(slide, x, Inches(1.7), col_w, col_h, fill=C_LGRAY)
        bullet_box(slide, col["items"],
                   x + Inches(0.15), Inches(1.78), col_w - Inches(0.3), col_h - Inches(0.15),
                   font_size=Pt(12), gap=5)


    # ── SLIDE 3: THE PROBLEM ────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "The challenge: viral signal in a sea of host RNA",
                 subtitle="Bulk RNA-seq is >99.95% human \u2014 k-mer classifiers with viral-only DBs create systematic false positives",
                 bar_color=C_CORAL)

    # Left: composition
    rect(slide, Inches(0.35), Inches(1.3), Inches(5.8), Inches(5.8), fill=C_LGRAY)
    add_text_box(slide, "Typical bulk RNA-seq library composition",
                 Inches(0.5), Inches(1.35), Inches(5.5), Inches(0.4),
                 font_size=Pt(13), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    bar_x, bar_y, bar_w = Inches(1.2), Inches(1.9), Inches(3.5)
    host_h = Inches(4.7)
    rect(slide, bar_x, bar_y, bar_w, host_h, fill=RGBColor(0xBB, 0xDE, 0xFB))
    add_text_box(slide, "Host (human)\n~99.95% of reads",
                 bar_x + Inches(0.1), bar_y + Inches(2.0), bar_w - Inches(0.2), Inches(0.7),
                 font_size=Pt(13), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    rect(slide, bar_x, bar_y, bar_w, Inches(0.22), fill=C_CORAL)
    add_text_box(slide, "Viral reads: <0.05%",
                 bar_x + Inches(0.1), bar_y + Inches(0.23), bar_w, Inches(0.35),
                 font_size=Pt(11), color=C_CORAL, bold=True)

    # Right: two failure modes
    rect(slide, Inches(6.4), Inches(1.3), Inches(6.55), Inches(5.8), fill=C_LGRAY)
    add_text_box(slide, "Two failure modes",
                 Inches(6.55), Inches(1.35), Inches(6.2), Inches(0.4),
                 font_size=Pt(13), bold=True, color=C_DARK)

    rect(slide, Inches(6.55), Inches(1.85), Inches(6.1), Inches(0.4), fill=C_CORAL)
    add_text_box(slide, "Problem 1: Viral-only database \u2014 the closed-world assumption",
                 Inches(6.65), Inches(1.87), Inches(6.0), Inches(0.38),
                 font_size=Pt(12), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    bullet_box(slide,
        ["No non-viral reference \u2014 every read must land somewhere viral",
         "Human transcripts sharing k-mers with viral genomes \u2192 misassigned",
         "HERVs, CMV k-mers, and neural lncRNAs are common offenders",
         "Result: 100% false positive rate in our baseline cohort",
        ],
        Inches(6.55), Inches(2.3), Inches(6.1), Inches(1.85),
        font_size=Pt(12), gap=5)

    rect(slide, Inches(6.55), Inches(4.3), Inches(6.1), Inches(0.4), fill=C_AMBER)
    add_text_box(slide, "Problem 2: Sensitivity \u2014 latent viruses leave few traces",
                 Inches(6.65), Inches(4.32), Inches(6.0), Inches(0.38),
                 font_size=Pt(12), bold=True, color=C_DARK)
    bullet_box(slide,
        ["VZV/HSV maintain latency with minimal transcription (LAT only)",
         "Empirical detection floor: ~10 reads per taxon at 60\u201375M read depth",
         "Below this, cannot distinguish signal from stochastic noise",
        ],
        Inches(6.55), Inches(4.75), Inches(6.1), Inches(1.5),
        font_size=Pt(12), gap=5)


    # ======================================================================
    # SECTION 2: HOW
    # ======================================================================
    section_divider(prs, 2,
        "How",
        "virome-pipeline: competitive dual-database classification from existing RNA-seq",
        color=C_BLUE)


    # ── SLIDE 4: PIPELINE ARCHITECTURE ──────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Pipeline architecture \u2014 v1.5.0",
                 subtitle="7-step Nextflow DSL2 workflow + dual-database branch + BLAST verification offshoot",
                 bar_color=C_BLUE)
    add_figure(slide, FIG1, x=Inches(0.35), y=Inches(1.25), w=Inches(12.6))

    # If no figure, draw simplified ASCII-style pipeline
    if not os.path.exists(FIG1):
        rect(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.5), fill=C_LGRAY)
        add_text_box(slide,
            "FASTQs \u2192 FASTQC \u2192 TRIMMOMATIC \u2192 STAR (host removal) \u2192 KRAKEN2 (DB1: viral) \u2192 BRACKEN \u2192 FILTER \u2192 AGGREGATE\n"
            "                                                                        \u2514\u2192 KRAKEN2 (DB2: PlusPF) \u2192 BRACKEN \u2192 FILTER \u2192 AGGREGATE \u2192 COMPARE_DB \u2192 tier classification\n\n"
            "BLAST offshoot (post-hoc, Tier 1 candidates):\n"
            "  [kraken2.output + unmapped FASTQs] \u2192 EXTRACT_READS \u2192 BLAST_VERIFY \u2192 ANALYZE \u2192 lifecycle report",
            Inches(0.7), Inches(1.8), Inches(11.9), Inches(3.5),
            font_name=CODE_FONT, font_size=Pt(12), color=C_DARK)
    footnote(slide, "All steps containerized (Apptainer) | SLURM-ready | Each process in its own .sif container | stageInMode = 'copy' for Apptainer symlink compatibility")


    # ── SLIDE 5: DUAL-DATABASE TIER SYSTEM ──────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "The core innovation: dual-database competitive classification",
                 subtitle="Run every sample through two Kraken2 databases simultaneously \u2014 only signals surviving both are real",
                 bar_color=C_BLUE)

    # Left: DB1
    rect(slide, Inches(0.35), Inches(1.25), Inches(3.8), Inches(5.75), fill=C_LGRAY)
    rect(slide, Inches(0.35), Inches(1.25), Inches(3.8), Inches(0.5), fill=C_CORAL)
    add_text_box(slide, "Database 1: Viral-only",
                 Inches(0.45), Inches(1.27), Inches(3.6), Inches(0.45),
                 font_size=Pt(14), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER)
    bullet_box(slide,
        ["~10,000 viral genomes",
         "No non-viral sequences",
         "Closed-world: every read forced into a viral bin",
         "Maximizes sensitivity",
         "But: 100% FP rate in baseline cohort",
        ],
        Inches(0.45), Inches(1.85), Inches(3.6), Inches(3.5),
        font_size=Pt(12.5), gap=6)

    # Middle: tiers
    rect(slide, Inches(4.4), Inches(1.25), Inches(4.5), Inches(5.75), fill=C_DARK)
    add_text_box(slide, "Three-tier confidence scoring",
                 Inches(4.5), Inches(1.32), Inches(4.3), Inches(0.45),
                 font_size=Pt(14), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER)

    tiers = [
        (C_GREEN, "Tier 1 \u2014 Shared",
         "Detected in BOTH databases\nHigh-confidence viral candidate\nUse for biology",
         "HSV-1 in PD19"),
        (C_AMBER, "Tier 2 \u2014 Viral-only",
         "Viral-only DB exclusive\nFalse positive candidates\nInvestigate before interpreting",
         "HERV-K, CMV proxy, MCV"),
        (C_MGRAY, "Tier 3 \u2014 PlusPF only",
         "PlusPF-only\nNon-viral contaminants\nQC utility \u2014 not biology",
         "10\u2013826 taxa / sample"),
    ]
    for i, (col, title, desc, result) in enumerate(tiers):
        y = Inches(1.95 + i * 1.65)
        rect(slide, Inches(4.55), y, Inches(4.2), Inches(1.45), fill=col)
        add_text_box(slide, title,
                     Inches(4.65), y + Inches(0.08), Inches(2.5), Inches(0.38),
                     font_size=Pt(13), bold=True, color=C_DARK)
        add_text_box(slide, desc,
                     Inches(4.65), y + Inches(0.48), Inches(2.5), Inches(0.9),
                     font_size=Pt(10.5), color=C_DARK)
        add_text_box(slide, result,
                     Inches(7.0), y + Inches(0.08), Inches(1.6), Inches(0.9),
                     font_size=Pt(11), bold=True, color=C_DARK, align=PP_ALIGN.RIGHT)

    # Right: DB2
    rect(slide, Inches(9.15), Inches(1.25), Inches(3.8), Inches(5.75), fill=C_LGRAY)
    rect(slide, Inches(9.15), Inches(1.25), Inches(3.8), Inches(0.5), fill=C_GREEN)
    add_text_box(slide, "Database 2: PlusPF (standard)",
                 Inches(9.25), Inches(1.27), Inches(3.6), Inches(0.45),
                 font_size=Pt(14), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER)
    bullet_box(slide,
        ["Full bacterial, archaeal, viral, fungal + human",
         "~100 GB \u2014 human genome included",
         "Open-world: reads compete against host",
         "HERV-K \u2192 reclassified to Homo sapiens",
         "CMV proxy \u2192 0 reads (artifact resolved)",
        ],
        Inches(9.25), Inches(1.85), Inches(3.6), Inches(3.5),
        font_size=Pt(12.5), gap=6)


    # ── SLIDE 6: MULTI-STAGE FILTERING ──────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Multi-stage filtering + curated artifact exclusion",
                 subtitle="Bracken raw \u2192 min-reads threshold (\u22655) \u2192 24-entry artifact exclusion list",
                 bar_color=C_BLUE)
    add_figure(slide, FIG2, x=Inches(0.6), y=Inches(1.25), w=Inches(5.5))

    # Right side: filtering explanation
    rect(slide, Inches(6.5), Inches(1.3), Inches(6.5), Inches(5.5), fill=C_LGRAY)
    add_text_box(slide, "Three output matrices per cohort",
                 Inches(6.65), Inches(1.38), Inches(6.2), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_DARK)

    stages = [
        (C_AMBER, "bracken_raw", "All viral species from Bracken\nNo threshold \u2014 baseline for QC"),
        (C_BLUE, "minreads", "\u22655 reads per taxon per sample\nValidated: noise floor ~10 reads at 60\u201375M depth"),
        (C_GREEN, "final (filtered)", "Curated artifact exclusion applied\n24 taxa: phages, reagent contaminants, env. metagenome, cross-mapping"),
    ]
    for i, (col, name, desc) in enumerate(stages):
        y = Inches(1.85 + i * 1.4)
        rect(slide, Inches(6.65), y, Inches(6.1), Inches(1.3), fill=RGBColor(0xFF, 0xFF, 0xFF))
        rect(slide, Inches(6.65), y, Inches(0.15), Inches(1.3), fill=col)
        add_text_box(slide, name,
                     Inches(6.9), y + Inches(0.1), Inches(2.5), Inches(0.35),
                     font_size=Pt(13), bold=True, color=col)
        add_text_box(slide, desc,
                     Inches(6.9), y + Inches(0.5), Inches(5.7), Inches(0.75),
                     font_size=Pt(11), color=C_DARK)

    add_text_box(slide, "Artifact categories",
                 Inches(6.65), Inches(6.15), Inches(6.2), Inches(0.3),
                 font_size=Pt(11), bold=True, color=C_DARK)
    add_text_box(slide,
        "Reagent contaminants \u00b7 ruminant orthobunyaviruses \u00b7 insect baculoviruses \u00b7 phages (ICTV-reclassified) \u00b7 "
        "DRG-specific env. metagenome cross-mapping \u00b7 avian herpesviruses \u00b7 giant amoeba viruses \u00b7 hantaviruses (Oxbow)",
        Inches(6.65), Inches(6.45), Inches(6.2), Inches(0.55),
        font_size=Pt(9.5), color=C_MGRAY)


    # ======================================================================
    # SECTION 3: RESULTS
    # ======================================================================
    section_divider(prs, 3,
        "Results",
        "Two cohorts, 36 donors, and the first confirmed viral detection in human DRG",
        color=C_CORAL)


    # ── SLIDE 7: COHORT OVERVIEW ─────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Study cohorts \u2014 36 unique donors across two runs",
                 subtitle="Non-PD baseline (n=16 unique) + Parkinson's DRG cohort (n=20)",
                 bar_color=C_CORAL)

    # Cohort 1: Non-PD baseline
    rect(slide, Inches(0.3), Inches(1.25), Inches(6.3), Inches(5.7), fill=C_LGRAY)
    rect(slide, Inches(0.3), Inches(1.25), Inches(6.3), Inches(0.55), fill=C_BLUE)
    add_text_box(slide, "Non-PD Baseline (v1.3.0)",
                 Inches(0.4), Inches(1.27), Inches(4.0), Inches(0.5),
                 font_size=Pt(15), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide, "n = 16 unique",
                 Inches(4.5), Inches(1.27), Inches(2.0), Inches(0.5),
                 font_size=Pt(18), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.RIGHT)
    bullet_box(slide,
        ["Skeletal muscle (n=5): independent donors, non-neural control",
         "DRG \u2014 donor1 (n=6): single donor, 6 spinal levels (L1\u2013L5, T12)",
         "DRG \u2014 Saad cohort (n=4): independent donors, published cohort",
         "AIG1390 (n=5): excluded \u2014 confirmed donor1 duplicate by MD5",
         "Saad_2: excluded \u2014 library failure (0 final taxa)",
         "",
         "Key result: Tier 1 = 0 across ALL samples",
         "100% false positive rate for viral-only Kraken2",
         "Established validated null baseline for DRG virome",
        ],
        Inches(0.4), Inches(1.9), Inches(6.0), Inches(4.5),
        font_size=Pt(12), gap=5)

    # Cohort 2: Parkinson's
    rect(slide, Inches(6.85), Inches(1.25), Inches(6.15), Inches(5.7), fill=C_LGRAY)
    rect(slide, Inches(6.85), Inches(1.25), Inches(6.15), Inches(0.55), fill=C_CORAL)
    add_text_box(slide, "Parkinson's DRG (v1.5.0)",
                 Inches(6.95), Inches(1.27), Inches(4.0), Inches(0.5),
                 font_size=Pt(15), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide, "n = 20",
                 Inches(10.8), Inches(1.27), Inches(2.1), Inches(0.5),
                 font_size=Pt(18), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.RIGHT)
    bullet_box(slide,
        ["PD patients (n=14): PD2\u2013PD6, PD9\u2013PD10, PD14\u2013PD20",
         "Unclassified controls (n=6): samples 023\u2013028",
         "All DRG tissue \u2014 Psomagen order AN00028264",
         "",
         "Key result: Tier 1 = 1",
         "HSV-1 detected in PD19 \u2014 46 reads, 1.89 RPM",
         "Confirmed in BOTH viral-only and PlusPF databases",
         "First Tier 1 detection across all 36 donors",
         "Absent from ALL other samples (0/35)",
        ],
        Inches(6.95), Inches(1.9), Inches(5.8), Inches(4.5),
        font_size=Pt(12), gap=5)


    # ── SLIDE 8: TIER RESULTS OVERVIEW ───────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Tier results across both cohorts",
                 subtitle="36 donors processed \u00b7 1 Tier 1 detection \u00b7 3 Tier 2 taxa fully resolved as false positives",
                 bar_color=C_CORAL)

    # Stat boxes
    stat_box(slide, "1",     "Tier 1 detection\n(HSV-1, PD19 only)",     Inches(0.4),  Inches(1.3), Inches(2.4), Inches(1.5), val_color=C_GREEN)
    stat_box(slide, "3",     "Tier 2 taxa\n(all resolved as FP)",        Inches(2.95), Inches(1.3), Inches(2.4), Inches(1.5), val_color=C_AMBER)
    stat_box(slide, "100%",  "Viral-only FP rate\n(baseline cohort)",    Inches(5.5),  Inches(1.3), Inches(2.4), Inches(1.5), val_color=C_CORAL)
    stat_box(slide, "24",    "Curated artifact\ntaxa excluded",          Inches(8.05), Inches(1.3), Inches(2.4), Inches(1.5), val_color=C_BLUE)
    stat_box(slide, "~10",   "RPM detection\nfloor (empirical)",         Inches(10.6), Inches(1.3), Inches(2.4), Inches(1.5), val_color=C_MGRAY)

    # Summary table
    headers = ["", "Tier 1 (confirmed)", "Tier 2 (FP candidates)", "Tier 3 (non-viral)"]
    col_ws = [Inches(3.0), Inches(3.1), Inches(3.1), Inches(3.5)]
    x_starts = []
    cx = Inches(0.35)
    for cw in col_ws:
        x_starts.append(cx)
        cx += cw + Inches(0.05)

    hdr_y = Inches(3.1)
    for j, (hdr, cw) in enumerate(zip(headers, col_ws)):
        rect(slide, x_starts[j], hdr_y, cw, Inches(0.42), fill=C_DARK)
        add_text_box(slide, hdr, x_starts[j] + Inches(0.05), hdr_y + Inches(0.04),
                     cw - Inches(0.1), Inches(0.35),
                     font_size=Pt(11), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    rows = [
        ("Non-PD baseline (n=16)", "0", "3: HERV-K, CMV proxy, MCV", "10\u2013179 taxa/sample"),
        ("Parkinson's DRG (n=20)", "1: HSV-1 in PD19 (46 reads)", "3: HERV-K, CMV proxy, MCV", "15\u2013826 taxa/sample"),
    ]
    for i, row in enumerate(rows):
        bg = C_LGRAY if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        row_y = Inches(3.55 + i * 0.72)
        for j, (cell, cw) in enumerate(zip(row, col_ws)):
            fill = bg
            if j == 1 and i == 1:
                fill = RGBColor(0xD5, 0xF5, 0xE3)  # green highlight for HSV-1
            rect(slide, x_starts[j], row_y, cw, Inches(0.68), fill=fill)
            add_text_box(slide, cell,
                         x_starts[j] + Inches(0.07), row_y + Inches(0.1),
                         cw - Inches(0.14), Inches(0.55),
                         font_size=Pt(11), color=C_DARK, word_wrap=True,
                         bold=(j == 0))

    # Key takeaway
    rect(slide, Inches(0.35), Inches(5.2), Inches(12.6), Inches(1.7), fill=C_DARK)
    add_text_box(slide, "Key interpretation",
                 Inches(0.5), Inches(5.28), Inches(12.3), Inches(0.35),
                 font_size=Pt(14), bold=True, color=C_AMBER)
    bullet_box(slide,
        ["The null result in the baseline cohort validates the pipeline \u2014 zero Tier 1 is the correct answer for these healthy donors",
         "HSV-1 in PD19 is the first Tier 1 detection \u2014 biologically plausible (HSV-1 latency in sensory ganglia), but n=1/14 PD (not statistically significant)",
         "All three Tier 2 taxa have fully characterized mechanisms: endogenous transcription (HERV-K), k-mer cross-mapping (CMV proxy), index hopping (MCV)",
        ],
        Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.2),
        font_size=Pt(11.5), color=RGBColor(0xEE, 0xEE, 0xEE), gap=4)


    # ── SLIDE 9: HSV-1 IN PD19 ──────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Tier 1: HSV-1 in PD19 \u2014 the first confirmed detection",
                 subtitle="46 reads, 1.89 RPM \u00b7 detected in both viral-only and PlusPF databases \u00b7 absent from all 35 other samples",
                 bar_color=C_GREEN)

    # Left: the data
    rect(slide, Inches(0.35), Inches(1.25), Inches(6.0), Inches(5.75), fill=C_LGRAY)
    rect(slide, Inches(0.35), Inches(1.25), Inches(6.0), Inches(0.5), fill=C_GREEN)
    add_text_box(slide, "The detection",
                 Inches(0.45), Inches(1.27), Inches(5.8), Inches(0.45),
                 font_size=Pt(15), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    stat_box(slide, "46",    "direct reads\n(Kraken2)",     Inches(0.5),  Inches(1.9), Inches(1.7), Inches(1.1), val_color=C_GREEN)
    stat_box(slide, "1.89",  "reads per\nmillion (RPM)",    Inches(2.3),  Inches(1.9), Inches(1.7), Inches(1.1), val_color=C_GREEN)
    stat_box(slide, "1/36",  "samples with\nthis detection", Inches(4.1), Inches(1.9), Inches(1.7), Inches(1.1), val_color=C_CORAL)

    add_text_box(slide, "Biological context",
                 Inches(0.5), Inches(3.2), Inches(5.7), Inches(0.35),
                 font_size=Pt(13), bold=True, color=C_DARK)
    bullet_box(slide,
        ["HSV-1 naturally establishes latency in sensory ganglia (DRG, TG)",
         "Latency marker: LAT (latency-associated transcript) \u2014 only RNA during dormancy",
         "Reactivation markers: immediate-early/early/late genes",
         "Epidemiological precedent: antiviral therapy associated with reduced PD risk",
         "Three proposed mechanisms:",
         ("Retrograde transport DRG \u2192 spinal cord \u2192 CNS", 1),
         ("\u03b1-synuclein aggregation triggered by viral proteins", 1),
         ("Neuroinflammation via STING/cGAS pathway", 1),
        ],
        Inches(0.5), Inches(3.6), Inches(5.7), Inches(3.0),
        font_size=Pt(11.5), gap=4)

    # Right: validation status
    rect(slide, Inches(6.6), Inches(1.25), Inches(6.4), Inches(5.75), fill=C_DARK)
    add_text_box(slide, "Validation status",
                 Inches(6.75), Inches(1.33), Inches(6.1), Inches(0.4),
                 font_size=Pt(15), bold=True, color=C_AMBER)

    add_text_box(slide, "Completed",
                 Inches(6.75), Inches(1.85), Inches(6.1), Inches(0.3),
                 font_size=Pt(12), bold=True, color=C_GREEN)
    bullet_box(slide,
        ["Dual-database confirmation: detected in BOTH DBs",
         "Tier 1 classification: high-confidence viral candidate",
         "Absent from all 35 other donors (disease specificity)",
         "Absent from non-PD baseline (n=16 unique samples)",
        ],
        Inches(6.75), Inches(2.2), Inches(6.1), Inches(1.8),
        font_size=Pt(11.5), color=RGBColor(0xEE, 0xEE, 0xEE), gap=5)

    add_text_box(slide, "Pending",
                 Inches(6.75), Inches(3.9), Inches(6.1), Inches(0.3),
                 font_size=Pt(12), bold=True, color=C_AMBER)
    bullet_box(slide,
        ["BLAST verification (blast_verify.nf ready \u2014 awaiting run)",
         "Life cycle phase: LAT (latency) vs. IE/E/L (reactivation)?",
         "Coverage depth: uniform genome coverage or LAT-localized?",
         "Independent validation (PCR, ISH) on remaining tissue",
        ],
        Inches(6.75), Inches(4.25), Inches(6.1), Inches(1.8),
        font_size=Pt(11.5), color=RGBColor(0xEE, 0xEE, 0xEE), gap=5)

    # Caveat box
    rect(slide, Inches(6.75), Inches(6.0), Inches(6.1), Inches(0.85), fill=C_CORAL)
    add_text_box(slide,
        "Statistical caveat: n=1/14 PD vs. 0/26 non-PD \u2014 not statistically significant (Fisher's p \u2248 0.35).\n"
        "This is a hypothesis-generating finding, not a conclusion. BLAST validation required.",
        Inches(6.85), Inches(6.05), Inches(5.9), Inches(0.75),
        font_size=Pt(10.5), color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)


    # ── SLIDE 10: TIER 2 FALSE POSITIVES ─────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Tier 2: the three false positives \u2014 fully resolved",
                 subtitle="Each illustrates a different class of artifact in human tissue metagenomics",
                 bar_color=C_AMBER)

    fps = [
        {
            "name": "HERV-K (HML-2)\ntaxon 45617",
            "reads": "Cohort total: ~80K reads",
            "color": C_PURPLE,
            "mech": "Endogenous retroviral transcription",
            "points": [
                "Present in ALL 36 samples (100% prevalence)",
                "5.8\u00d7 enriched in DRG vs. muscle (p = 3.3\u00d710\u207b\u2074)",
                "DRG mean: 42 RPM (non-PD), 32 RPM (PD)",
                "PlusPF: ALL reads \u2192 Homo sapiens",
                "Neural-lineage LTR transcription in sensory neurons",
                "Real biology, but endogenous \u2014 not infection",
            ]
        },
        {
            "name": "CMV proxy\ntaxon 3050337",
            "reads": "~1,100 reads across 28 samples",
            "color": C_CORAL,
            "mech": "Database gap + k-mer cross-mapping",
            "points": [
                "Cytomegalovirus papiinebeta3 (baboon CMV)",
                "ICTV 2023 split HHV-5 into primate species",
                "Human reads routed via baboon CMV child node",
                "~10\u00d7 DRG-enriched (same k-mer source as HERV-K)",
                "PlusPF: ZERO reads \u2014 complete resolution",
                "Relabeled \u2192 'Human CMV (HHV-5) [proxy]'",
            ]
        },
        {
            "name": "Molluscum contagiosum\ntaxon 10279",
            "reads": "~190 reads, sporadic",
            "color": C_AMBER,
            "mech": "Index hopping / contamination",
            "points": [
                "Sporadic, tissue-independent distribution",
                "Absent from deepest library (Saad_1, 430K reads)",
                "Muscle-predominant (65%) \u2014 not DRG-specific",
                "Higher prevalence in PD cohort (11/14 vs. 24%)",
                "PlusPF: not detected \u2014 no genuine signal",
                "Consistent with patterned flow cell index hopping",
            ]
        },
    ]
    for i, fp in enumerate(fps):
        x = Inches(0.3 + i * 4.35)
        rect(slide, x, Inches(1.3), Inches(4.15), Inches(5.7), fill=C_LGRAY)
        rect(slide, x, Inches(1.3), Inches(4.15), Inches(0.72), fill=fp["color"])
        add_text_box(slide, fp["name"],
                     x + Inches(0.1), Inches(1.3), Inches(3.0), Inches(0.7),
                     font_size=Pt(12), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        add_text_box(slide, fp["reads"],
                     x + Inches(0.1), Inches(2.08), Inches(3.9), Inches(0.3),
                     font_size=Pt(10.5), color=C_DARK, bold=True, italic=True)
        add_text_box(slide, fp["mech"],
                     x + Inches(0.1), Inches(2.38), Inches(3.9), Inches(0.3),
                     font_size=Pt(10.5), color=fp["color"], bold=True)
        bullet_box(slide, fp["points"],
                   x + Inches(0.1), Inches(2.72), Inches(3.9), Inches(3.3),
                   font_size=Pt(11), gap=5)


    # ── SLIDE 11: HERV-K DEEP DIVE ──────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "HERV-K \u2014 a real biological signal worth tracking",
                 subtitle="5.8\u00d7 DRG enrichment over muscle \u00b7 reproduced across cohorts \u00b7 first quantification of HERV-K in human DRG",
                 bar_color=C_PURPLE)
    add_figure(slide, FIG4, x=Inches(0.35), y=Inches(1.25), w=Inches(7.0))

    # Right panel
    rect(slide, Inches(7.6), Inches(1.3), Inches(5.4), Inches(5.7), fill=C_LGRAY)
    add_text_box(slide, "Cross-cohort HERV-K comparison",
                 Inches(7.75), Inches(1.38), Inches(5.1), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_DARK)

    # Mini table
    rows = [
        ("Muscle (non-PD)", "n=5", "31.3 RPM"),
        ("DRG \u2014 donor1", "n=6", "45.8 RPM"),
        ("DRG \u2014 Saad", "n=4", "72.7 RPM"),
        ("DRG \u2014 PD patients", "n=14", "31.8 RPM"),
        ("DRG \u2014 PD controls", "n=6", "34.5 RPM"),
    ]
    for i, (group, n, rpm) in enumerate(rows):
        y = Inches(1.9 + i * 0.45)
        bg = C_STRIPE if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        rect(slide, Inches(7.75), y, Inches(5.1), Inches(0.42), fill=bg)
        add_text_box(slide, group, Inches(7.85), y + Inches(0.06), Inches(2.5), Inches(0.3),
                     font_size=Pt(10.5), color=C_DARK)
        add_text_box(slide, n, Inches(10.4), y + Inches(0.06), Inches(0.8), Inches(0.3),
                     font_size=Pt(10.5), color=C_MGRAY, align=PP_ALIGN.CENTER)
        add_text_box(slide, rpm, Inches(11.3), y + Inches(0.06), Inches(1.4), Inches(0.3),
                     font_size=Pt(10.5), bold=True, color=C_PURPLE, align=PP_ALIGN.RIGHT)

    add_text_box(slide, "Implications",
                 Inches(7.75), Inches(4.3), Inches(5.1), Inches(0.3),
                 font_size=Pt(13), bold=True, color=C_DARK)
    bullet_box(slide,
        ["HERV-K is NOT a virus \u2014 it's endogenous retroviral chromosomal transcription",
         "DRG enrichment likely reflects neural-lineage chromatin accessibility",
         "ALS literature: HERV-K Env protein in motor neurons (Li et al. 2015)",
         "PD de-repression hypothesis NOT supported (PD RPM \u2264 non-PD RPM)",
         "Could serve as internal control \u2014 HERV-K RPM as library quality metric",
        ],
        Inches(7.75), Inches(4.65), Inches(5.1), Inches(2.3),
        font_size=Pt(11), gap=5)


    # ── SLIDE 12: DIVERSITY COMPARISON ───────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Viral diversity across cohorts \u2014 sparse and dominated by endogenous signal",
                 subtitle="Mean richness 2\u20134 taxa per sample \u00b7 Shannon diversity <0.5 \u00b7 HERV-K accounts for >90% of viral reads",
                 bar_color=C_CORAL)

    # Summary stats by group
    groups = [
        ("Non-PD Muscle", "n=5", "2.0", "0.067", "~630", C_BLUE),
        ("Non-PD DRG (donor1)", "n=6", "3.0", "0.185", "~3,670", C_PURPLE),
        ("Non-PD DRG (Saad)", "n=4", "2.8", "0.390", "~5,800", C_GREEN),
        ("PD DRG patients", "n=14", "3.1", "0.280", "~800", C_CORAL),
        ("PD controls (023\u2013028)", "n=6", "2.8", "0.226", "~750", C_AMBER),
    ]

    headers = ["Group", "n", "Richness", "Shannon H'", "Mean reads", ""]
    col_ws = [Inches(3.0), Inches(0.8), Inches(1.4), Inches(1.4), Inches(1.6), Inches(4.5)]
    x_starts = []
    cx = Inches(0.35)
    for cw in col_ws:
        x_starts.append(cx)
        cx += cw + Inches(0.05)

    # Headers
    for j, (hdr, cw) in enumerate(zip(headers, col_ws)):
        if j < 5:
            rect(slide, x_starts[j], Inches(1.3), cw, Inches(0.42), fill=C_DARK)
            add_text_box(slide, hdr, x_starts[j] + Inches(0.05), Inches(1.34),
                         cw - Inches(0.1), Inches(0.35),
                         font_size=Pt(11), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    for i, (group, n, rich, shannon, reads, col) in enumerate(groups):
        bg = C_LGRAY if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        y = Inches(1.75 + i * 0.55)
        row_data = [group, n, rich, shannon, reads]
        for j, (cell, cw) in enumerate(zip(row_data, col_ws)):
            rect(slide, x_starts[j], y, cw, Inches(0.5), fill=bg)
            rect(slide, x_starts[j], y, Inches(0.12) if j == 0 else 0, Inches(0.5), fill=col)
            add_text_box(slide, cell,
                         x_starts[j] + Inches(0.15), y + Inches(0.08),
                         cw - Inches(0.2), Inches(0.35),
                         font_size=Pt(11), color=C_DARK,
                         bold=(j == 0), align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    # Right side: interpretation
    rect(slide, Inches(7.65), Inches(1.3), Inches(5.35), Inches(5.7), fill=C_LGRAY)
    add_text_box(slide, "What low diversity means",
                 Inches(7.8), Inches(1.38), Inches(5.0), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_DARK)
    bullet_box(slide,
        ["The DRG virome (from bulk RNA-seq) is extremely sparse",
         "HERV-K dominates: >90% of total viral reads in every sample",
         "After artifact exclusion, 2\u20134 taxa per sample \u2014 most are Tier 2",
         "This is expected: bulk RNA-seq is not designed for metagenomics",
         "Latent viruses (VZV, HSV) may be below detection floor in most donors",
         "",
         "HSV-1 in PD19 is the outlier \u2014 46 reads is well above the ~10 read noise floor, "
         "and it survived competitive classification against the full genome",
        ],
        Inches(7.8), Inches(1.85), Inches(5.0), Inches(4.5),
        font_size=Pt(12), gap=6)


    # ======================================================================
    # SECTION 4: FUTURE DIRECTIONS
    # ======================================================================
    section_divider(prs, 4,
        "Future Directions",
        "Compiling the lab's data \u2014 beginning with bulk RNA-seq",
        color=C_GREEN)


    # ── SLIDE 13: IMMEDIATE PRIORITIES ───────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Immediate next steps",
                 subtitle="BLAST validation \u00b7 Paper 1 submission \u00b7 begin lab-wide data compilation",
                 bar_color=C_GREEN)

    # Column 1: Validate
    rect(slide, Inches(0.3), Inches(1.25), Inches(4.15), Inches(5.75), fill=C_LGRAY)
    rect(slide, Inches(0.3), Inches(1.25), Inches(4.15), Inches(0.5), fill=C_CORAL)
    add_text_box(slide, "Validate PD19 HSV-1",
                 Inches(0.4), Inches(1.27), Inches(3.9), Inches(0.45),
                 font_size=Pt(14), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    bullet_box(slide,
        ["Run blast_verify.nf on PD19 (pipeline ready, awaiting execution)",
         "Determine life cycle phase: LAT (latency) vs. IE/E/L (reactivation)",
         "Map reads to HSV-1 reference genome (NC_001806.2) for coverage profile",
         "Compare to known LAT locus coordinates",
         "If confirmed: first RNA-seq evidence of HSV-1 in Parkinson's DRG",
        ],
        Inches(0.4), Inches(1.85), Inches(3.9), Inches(4.5),
        font_size=Pt(12), gap=7)

    # Column 2: Publish
    rect(slide, Inches(4.7), Inches(1.25), Inches(4.0), Inches(5.75), fill=C_LGRAY)
    rect(slide, Inches(4.7), Inches(1.25), Inches(4.0), Inches(0.5), fill=C_BLUE)
    add_text_box(slide, "Paper 1: Baseline",
                 Inches(4.8), Inches(1.27), Inches(3.8), Inches(0.45),
                 font_size=Pt(14), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    bullet_box(slide,
        ["Bioinformatics Application Note (~2,500 words)",
         "Pipeline method + null baseline result",
         "Dual-DB competitive classification validated",
         "100% FP elimination demonstrated",
         "HERV-K DRG enrichment characterized",
         "Zenodo DOI + SRA data submission",
         "Draft in progress \u2014 abstract written",
        ],
        Inches(4.8), Inches(1.85), Inches(3.8), Inches(4.5),
        font_size=Pt(12), gap=7)

    # Column 3: Compile
    rect(slide, Inches(8.95), Inches(1.25), Inches(4.05), Inches(5.75), fill=C_DARK)
    add_text_box(slide, "Compile lab bulk RNA-seq",
                 Inches(9.05), Inches(1.33), Inches(3.8), Inches(0.45),
                 font_size=Pt(14), bold=True, color=C_AMBER)
    bullet_box(slide,
        ["Inventory all existing bulk RNA-seq across the lab",
         "Prioritize: DRG, spinal cord, trigeminal ganglia",
         "Consolidate FASTQs to shared location on /groups/tprice/",
         "Build master samplesheet with metadata",
         "Run virome-pipeline on everything \u2014 one Nextflow command",
         "Goal: 100+ donor DRG virome atlas from existing data",
         "No new sequencing needed for first pass",
        ],
        Inches(9.05), Inches(1.85), Inches(3.8), Inches(4.5),
        font_size=Pt(12), color=RGBColor(0xEE, 0xEE, 0xEE), gap=7)


    # ── SLIDE 14: LAB DATA COMPILATION ───────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Compiling the lab's bulk RNA-seq \u2014 the path to an atlas",
                 subtitle="The pipeline is ready \u2014 the bottleneck is organizing the data we already have",
                 bar_color=C_GREEN)

    # Left: what we need
    rect(slide, Inches(0.3), Inches(1.25), Inches(6.3), Inches(5.75), fill=C_LGRAY)
    add_text_box(slide, "What's needed to begin",
                 Inches(0.45), Inches(1.33), Inches(6.0), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_DARK)
    bullet_box(slide,
        ["Inventory: which donors have bulk RNA-seq? Which tissues? Where are the FASTQs?",
         "Metadata: tissue type, disease status, donor demographics, sequencing batch",
         "Central storage: consolidate to /groups/tprice/ or /scratch/ with clear directory structure",
         "Samplesheet: CSV with sample, fastq_r1, fastq_r2 \u2014 pipeline handles the rest",
         "Tissue-specific artifact lists: DRG list exists (24 taxa); other tissues need curation",
         "Flag paired-tissue designs: DRG + muscle from same donor is especially powerful",
        ],
        Inches(0.45), Inches(1.8), Inches(6.0), Inches(4.5),
        font_size=Pt(12.5), gap=8)

    # Right: what the pipeline produces
    rect(slide, Inches(6.85), Inches(1.25), Inches(6.15), Inches(5.75), fill=C_DARK)
    add_text_box(slide, "What the pipeline produces per cohort",
                 Inches(7.0), Inches(1.33), Inches(5.8), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_AMBER)

    outputs = [
        ("viral_abundance_matrix.tsv", "Primary output \u2014 reads + RPM per taxon per sample"),
        ("consensus_matrix.tsv", "Tier 1 only \u2014 high-confidence viral candidates"),
        ("db_comparison_summary.tsv", "Per-sample tier counts (Tier 1/2/3 breakdown)"),
        ("virome_report/summary.html", "Interactive HTML: heatmap, funnel, diversity, tier plot"),
        ("multiqc_report.html", "QC: FastQC, STAR alignment, Bracken classification stats"),
        ("filter_summary.tsv", "Per-stage taxa/read counts for QC tracking"),
    ]
    for i, (fname, desc) in enumerate(outputs):
        y = Inches(1.85 + i * 0.78)
        rect(slide, Inches(7.0), y, Inches(5.8), Inches(0.72), fill=RGBColor(0x22, 0x22, 0x3E))
        add_text_box(slide, fname,
                     Inches(7.1), y + Inches(0.05), Inches(5.6), Inches(0.3),
                     font_name=CODE_FONT, font_size=Pt(10.5), color=C_AMBER, bold=True)
        add_text_box(slide, desc,
                     Inches(7.1), y + Inches(0.35), Inches(5.6), Inches(0.35),
                     font_size=Pt(10.5), color=RGBColor(0xCC, 0xCC, 0xCC))

    add_text_box(slide,
        "One Nextflow command per cohort. Results in hours, not weeks.",
        Inches(7.0), Inches(6.55), Inches(5.8), Inches(0.35),
        font_size=Pt(12), bold=True, color=C_GREEN)


    # ── SLIDE 15: LONGER-TERM ROADMAP ────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Longer-term roadmap \u2014 from bulk to single-cell to spatial",
                 subtitle="Each layer adds resolution; all leverage the same dual-database classification framework",
                 bar_color=C_GREEN)

    phases = [
        {
            "title": "Now \u2192 3 months",
            "color": C_CORAL,
            "items": [
                "BLAST validation of PD19 HSV-1",
                "Paper 1 submission (baseline \u2014 Bioinformatics)",
                "Lab bulk RNA-seq inventory + compilation",
                "100+ donor DRG atlas (first pass)",
                "Identify samples 023\u2013028 + missing PD IDs",
            ]
        },
        {
            "title": "3 \u2192 9 months",
            "color": C_BLUE,
            "items": [
                "Paper 2: PD vs. non-PD DRG virome (J. Virology)",
                "Tissue expansion: spinal cord, TG, vagus nerve",
                "Clinical metadata integration (diagnosis, age, PMI)",
                "DESeq2-style differential abundance testing",
                "Alignment-based validation module (minimap2/PathSeq)",
            ]
        },
        {
            "title": "9 \u2192 18 months",
            "color": C_GREEN,
            "items": [
                "Pseudobulk virome from snRNA-seq + Visium",
                "Cell-type resolved viral tropism (which cells harbor virus?)",
                "De novo assembly for divergent/novel viruses (SPAdes/MEGAHIT)",
                "Multi-cohort atlas: pain conditions vs. healthy donors",
                "Human neural virome atlas paper (Nature Methods / Cell Genomics)",
            ]
        },
    ]
    for i, phase in enumerate(phases):
        x = Inches(0.3 + i * 4.35)
        rect(slide, x, Inches(1.25), Inches(4.15), Inches(5.75), fill=C_LGRAY)
        rect(slide, x, Inches(1.25), Inches(4.15), Inches(0.55), fill=phase["color"])
        add_text_box(slide, phase["title"],
                     x + Inches(0.1), Inches(1.27), Inches(3.9), Inches(0.5),
                     font_size=Pt(15), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
        bullet_box(slide, phase["items"],
                   x + Inches(0.15), Inches(1.9), Inches(3.85), Inches(4.8),
                   font_size=Pt(12.5), gap=8)


    # ── SLIDE 16: CLOSING ────────────────────────────────────────────────
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, fill=C_DARK)
    rect(slide, 0, 0, W, Inches(0.1), fill=C_AMBER)
    rect(slide, 0, H - Inches(0.1), W, Inches(0.1), fill=C_AMBER)

    add_text_box(slide, "Summary",
                 Inches(0.6), Inches(0.5), Inches(4.0), Inches(0.5),
                 font_size=Pt(28), bold=True, color=C_AMBER)

    summary_pts = [
        "The DRG is a known viral latency reservoir \u2014 but no one has systematically profiled its virome from RNA-seq",
        "virome-pipeline solves the closed-world assumption via competitive dual-database classification",
        "Baseline cohort (n=16): zero Tier 1 \u2014 validated null result; 100% FP rate for viral-only DB eliminated",
        "Parkinson's cohort (n=20): first Tier 1 detection \u2014 HSV-1 in PD19 (46 reads, 1.89 RPM)",
        "Three Tier 2 false positives fully resolved: HERV-K (endogenous), CMV proxy (k-mer artifact), MCV (index hopping)",
        "HERV-K 5.8\u00d7 DRG-enriched \u2014 first quantification in human DRG; biologically real but not infection",
        "The pipeline is ready to scale \u2014 the next step is compiling the lab's bulk RNA-seq data",
    ]
    bullet_box(slide, summary_pts,
               Inches(0.6), Inches(1.1), Inches(12.0), Inches(4.0),
               font_size=Pt(13), color=RGBColor(0xEE, 0xEE, 0xEE), gap=7)

    add_text_box(slide, "The vision",
                 Inches(0.6), Inches(5.2), Inches(12.0), Inches(0.45),
                 font_size=Pt(18), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(slide,
        "A comprehensive human neural virome atlas \u2014 100+ donors, multiple tissue types and pain conditions \u2014 "
        "built from the sequencing data this lab already has.",
        Inches(0.6), Inches(5.7), Inches(12.0), Inches(0.7),
        font_size=Pt(15), color=C_AMBER, word_wrap=True)

    add_text_box(slide, "virome-pipeline v1.5.0  \u00b7  Nextflow DSL2  \u00b7  MIT License",
                 Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.4),
                 font_size=Pt(11), color=RGBColor(0x88, 0x88, 0x88))
    add_text_box(slide, "Questions?",
                 Inches(9.5), Inches(5.2), Inches(3.3), Inches(0.8),
                 font_size=Pt(36), bold=True, color=C_GREEN, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    prs = new_prs()
    build(prs)
    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")
