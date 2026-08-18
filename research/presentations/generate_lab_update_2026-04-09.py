"""
generate_lab_update_2026-04-09.py
Brief lab update: DRG Virome Pipeline — Parkinson's Cohort Results
2026-04-09 — Matthew Wild, TJP Lab, UT Dallas

Short, figure-heavy deck for a quick lab meeting update.

Run from repo root:
    python3 research/presentations/generate_lab_update_2026-04-09.py

Outputs:
    research/presentations/virome_lab_update_2026-04-09.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
REPO_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
OUT_PATH = os.path.join(
    REPO_ROOT, "research", "presentations",
    "virome_lab_update_2026-04-09.pptx",
)

HEATMAP = os.path.join(
    REPO_ROOT, "results", "parkinson_2026", "results",
    "virome_report", "heatmap.png",
)
DB_COMPARISON = os.path.join(
    REPO_ROOT, "results", "parkinson_2026", "results",
    "db_comparison", "db_comparison.png",
)

# ---------------------------------------------------------------------------
# Design constants (16:9)
# ---------------------------------------------------------------------------
W = Inches(13.33)
H = Inches(7.5)

C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
C_BLUE   = RGBColor(0x2E, 0x86, 0xAB)
C_CORAL  = RGBColor(0xE8, 0x48, 0x55)
C_GREEN  = RGBColor(0x3B, 0xB2, 0x73)
C_AMBER  = RGBColor(0xF9, 0xC7, 0x4F)
C_PURPLE = RGBColor(0x7B, 0x2D, 0x8B)
C_LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
C_MGRAY  = RGBColor(0x9E, 0x9E, 0x9E)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(slide, x, y, w, h, fill=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    return s


def txt(slide, text, x, y, w, h, size=14, bold=False, italic=False,
        color=C_DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(slide, items, x, y, w, h, size=13, gap=6, color=C_DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        pPr = p._p.get_or_add_pPr()
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', str(gap * 100))
        r = p.add_run()
        r.text = "\u25b8 " + item
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def header(slide, title, bar_color=C_BLUE):
    box(slide, 0, 0, W, Inches(0.08), fill=bar_color)
    txt(slide, title, Inches(0.45), Inches(0.15), Inches(12.4), Inches(0.6),
        size=26, bold=True, color=C_DARK)
    box(slide, 0, Inches(0.82), W, Inches(0.02), fill=C_LGRAY)


def fig(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        box(slide, x, y, w, h or Inches(3), fill=C_LGRAY)
        txt(slide, f"[{os.path.basename(path)}]",
            x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), Inches(0.4),
            size=10, color=C_MGRAY)
        return
    if h:
        slide.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        slide.shapes.add_picture(path, x, y, width=w)


def footnote(slide, text):
    txt(slide, text, Inches(0.45), Inches(7.1), Inches(12.4), Inches(0.35),
        size=9, color=C_MGRAY, italic=True)


def stat(slide, value, label, x, y, w=Inches(2.2), h=Inches(1.1),
         val_color=C_BLUE):
    box(slide, x, y, w, h, fill=C_LGRAY)
    txt(slide, value, x, y + Inches(0.08), w, Inches(0.55),
        size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + Inches(0.62), w, Inches(0.4),
        size=10, color=C_DARK, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build(prs):

    # ── SLIDE 1: TITLE ──────────────────────────────────────────────────
    slide = blank(prs)
    box(slide, 0, 0, W, Inches(4.5), fill=C_DARK)
    box(slide, 0, Inches(4.5), W, Inches(3.0), fill=C_BLUE)

    txt(slide, "DRG Virome Pipeline Update",
        Inches(0.6), Inches(1.0), Inches(10.5), Inches(1.2),
        size=40, bold=True, color=C_WHITE)
    txt(slide, "Parkinson's DRG cohort (n=20) \u2014 first Tier 1 viral detection",
        Inches(0.6), Inches(2.3), Inches(10.5), Inches(0.6),
        size=18, color=C_AMBER)
    txt(slide, "Matthew Wild  \u00b7  TJP Lab  \u00b7  April 2026",
        Inches(0.6), Inches(3.1), Inches(10.5), Inches(0.5),
        size=14, color=RGBColor(0xBB, 0xBB, 0xBB))

    txt(slide, "v1.5.0",
        Inches(10.5), Inches(1.1), Inches(2.5), Inches(0.4),
        size=12, color=C_AMBER, bold=True, align=PP_ALIGN.RIGHT)

    # Bottom bar: agenda
    txt(slide, "Background  \u2192  Pipeline  \u2192  PD cohort results  \u2192  Next steps",
        Inches(0.6), Inches(5.2), Inches(12.0), Inches(0.5),
        size=16, bold=True, color=C_WHITE)


    # ── SLIDE 2: CONTEXT ────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Background", bar_color=C_PURPLE)

    # Left column
    txt(slide, "The question",
        Inches(0.45), Inches(1.0), Inches(5.5), Inches(0.4),
        size=16, bold=True, color=C_PURPLE)
    bullets(slide, [
        "Neurotropic viruses (HSV-1, VZV, CMV) establish latency in DRG",
        "Viral reactivation implicated in neuropathic pain, neurodegeneration",
        "HSV-1 linked to Parkinson's in epidemiological studies",
        "No systematic virome profiling of human DRG has been reported",
    ], Inches(0.45), Inches(1.5), Inches(5.5), Inches(3.0), size=13, gap=8)

    # Right column
    txt(slide, "The approach",
        Inches(7.0), Inches(1.0), Inches(5.5), Inches(0.4),
        size=16, bold=True, color=C_BLUE)
    bullets(slide, [
        "Use existing bulk RNA-seq \u2014 host-deplete with STAR, classify with Kraken2",
        "Dual-database classification: viral-only vs. PlusPF (full genome)",
        "Only taxa detected in BOTH databases = Tier 1 (confirmed)",
        "Baseline cohort (16 non-PD donors): Tier 1 = 0 \u2014 validated null",
    ], Inches(7.0), Inches(1.5), Inches(5.8), Inches(3.0), size=13, gap=8)

    # Bottom: what's new
    box(slide, Inches(0.3), Inches(4.8), Inches(12.7), Inches(2.3), fill=C_DARK)
    txt(slide, "What's new: Parkinson's DRG cohort",
        Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.4),
        size=16, bold=True, color=C_AMBER)
    bullets(slide, [
        "20 new DRG samples: 14 PD patients + 6 unclassified controls (Psomagen AN00028264)",
        "Same pipeline, same dual-database classification \u2014 first application to a disease cohort",
        "Result: the first Tier 1 detection across all 36 donors profiled to date",
    ], Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5),
       size=13, gap=7, color=RGBColor(0xEE, 0xEE, 0xEE))


    # ── SLIDE 3: HEATMAP ────────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Viral abundance heatmap \u2014 PD cohort (n=20)", bar_color=C_CORAL)

    fig(slide, HEATMAP, x=Inches(0.3), y=Inches(1.0), w=Inches(8.8))

    # Right annotation panel
    box(slide, Inches(9.4), Inches(1.0), Inches(3.6), Inches(5.8), fill=C_LGRAY)
    txt(slide, "Key observations",
        Inches(9.55), Inches(1.1), Inches(3.3), Inches(0.35),
        size=14, bold=True, color=C_DARK)
    bullets(slide, [
        "HERV-K dominant in all samples (endogenous, not infection)",
        "CMV proxy present but Tier 2 (k-mer artifact)",
        "MCV sporadic (index hopping)",
        "Simplexvirus humanalpha1 = HSV-1",
        "HSV-1 detected ONLY in PD19",
        "46 reads, 1.89 RPM",
        "Absent from all other 19 samples",
    ], Inches(9.55), Inches(1.55), Inches(3.3), Inches(4.5), size=11.5, gap=6)

    # Arrow callout pointing to HSV row
    box(slide, Inches(9.4), Inches(5.7), Inches(3.6), Inches(1.0), fill=C_GREEN)
    txt(slide, "HSV-1 in PD19 is the\nonly Tier 1 detection\nacross all 36 donors",
        Inches(9.5), Inches(5.75), Inches(3.4), Inches(0.9),
        size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    footnote(slide, "log10(reads + 1) scale  |  Final filtered matrix (post-artifact exclusion)  |  'Human CMV (HHV-5) [proxy]' = k-mer artifact, not genuine CMV")


    # ── SLIDE 4: DB COMPARISON ───────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Dual-database tier classification \u2014 PD cohort", bar_color=C_BLUE)

    # Figure takes most of the slide
    fig(slide, DB_COMPARISON,
        x=Inches(0.3), y=Inches(1.2), w=Inches(12.7), h=Inches(3.9))

    # Annotation below
    box(slide, Inches(0.3), Inches(5.3), Inches(12.7), Inches(2.0), fill=C_LGRAY)

    # Three tier labels
    tier_data = [
        (C_GREEN,  "Shared (Tier 1)",      "Detected in BOTH DBs = confirmed viral signal",    "1 taxon in 1 sample (HSV-1, PD19)"),
        (RGBColor(0xF0, 0x96, 0x46), "Viral-only (Tier 2)",  "Viral-only DB exclusive = false positive candidates", "3 taxa per sample (HERV-K, CMV proxy, MCV)"),
        (C_PURPLE, "PlusPF only (Tier 3)", "Non-viral background contaminants",                "15\u201330 taxa per sample (QC utility)"),
    ]
    for i, (col, name, desc, result) in enumerate(tier_data):
        x = Inches(0.5 + i * 4.3)
        box(slide, x, Inches(5.4), Inches(0.25), Inches(0.25), fill=col)
        txt(slide, name, x + Inches(0.35), Inches(5.37), Inches(2.5), Inches(0.3),
            size=12, bold=True, color=C_DARK)
        txt(slide, desc, x + Inches(0.35), Inches(5.68), Inches(3.8), Inches(0.3),
            size=10, color=C_MGRAY)
        txt(slide, result, x + Inches(0.35), Inches(5.98), Inches(3.8), Inches(0.3),
            size=11, color=col, bold=True)

    footnote(slide, "PD19 is the only sample with a green (Shared) segment  |  Baseline cohort (n=16 non-PD donors) had zero Tier 1 detections")


    # ── SLIDE 5: PD19 HSV-1 ─────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "PD19 HSV-1 \u2014 what we know and what's next", bar_color=C_GREEN)

    # Stats row
    stat(slide, "46",   "reads (Kraken2)",    Inches(0.4),  Inches(1.0), val_color=C_GREEN)
    stat(slide, "1.89", "RPM",                Inches(2.75), Inches(1.0), val_color=C_GREEN)
    stat(slide, "1/36", "donors with HSV-1",  Inches(5.1),  Inches(1.0), val_color=C_CORAL)
    stat(slide, "0/22", "non-PD with HSV-1",  Inches(7.45), Inches(1.0), val_color=C_BLUE)
    stat(slide, "1/14", "PD with HSV-1",      Inches(9.8),  Inches(1.0), val_color=C_PURPLE)

    # Two columns
    txt(slide, "Why this is plausible",
        Inches(0.45), Inches(2.4), Inches(5.5), Inches(0.4),
        size=15, bold=True, color=C_GREEN)
    bullets(slide, [
        "HSV-1 naturally establishes latency in sensory ganglia",
        "Antiviral therapy linked to reduced PD risk in registries",
        "Proposed mechanisms: retrograde transport, \u03b1-synuclein interaction, neuroinflammation",
        "46 reads well above ~10 read noise floor",
        "Survived competitive classification against full human genome",
    ], Inches(0.45), Inches(2.9), Inches(5.8), Inches(3.0), size=12.5, gap=7)

    txt(slide, "What's pending",
        Inches(7.0), Inches(2.4), Inches(5.5), Inches(0.4),
        size=15, bold=True, color=C_AMBER)
    bullets(slide, [
        "BLAST verification (blast_verify.nf ready to run)",
        "Life cycle phase: latency (LAT) or reactivation (IE/E/L)?",
        "Genome coverage profile against NC_001806.2",
        "Not statistically significant yet (Fisher's p \u2248 0.35, n=1)",
        "Hypothesis-generating \u2014 needs larger PD cohort",
    ], Inches(7.0), Inches(2.9), Inches(5.8), Inches(3.0), size=12.5, gap=7)

    # Bottom caveat
    box(slide, Inches(0.3), Inches(6.15), Inches(12.7), Inches(0.8), fill=C_DARK)
    txt(slide,
        "This is a hypothesis, not a conclusion. BLAST validation and a larger cohort are required before any biological claim.",
        Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5),
        size=13, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER)


    # ── SLIDE 6: NEXT STEPS ─────────────────────────────────────────────
    slide = blank(prs)
    header(slide, "Next steps", bar_color=C_GREEN)

    cols = [
        {
            "title": "This month",
            "color": C_CORAL,
            "items": [
                "BLAST validation of PD19 HSV-1",
                "Determine latency vs. reactivation",
                "Identify samples 023\u2013028",
            ]
        },
        {
            "title": "This quarter",
            "color": C_BLUE,
            "items": [
                "Paper 1 submission (baseline null result)",
                "Begin lab-wide bulk RNA-seq inventory",
                "Compile samplesheets for existing data",
            ]
        },
        {
            "title": "Longer term",
            "color": C_GREEN,
            "items": [
                "100+ donor DRG virome atlas",
                "PD vs. non-PD comparison paper",
                "Expand to snRNA-seq / spatial (pseudobulk)",
            ]
        },
    ]
    for i, col in enumerate(cols):
        x = Inches(0.3 + i * 4.35)
        box(slide, x, Inches(1.0), Inches(4.15), Inches(0.55), fill=col["color"])
        txt(slide, col["title"], x + Inches(0.15), Inches(1.03), Inches(3.85), Inches(0.5),
            size=16, bold=True, color=C_WHITE)
        box(slide, x, Inches(1.55), Inches(4.15), Inches(3.8), fill=C_LGRAY)
        bullets(slide, col["items"],
                x + Inches(0.15), Inches(1.7), Inches(3.85), Inches(3.5),
                size=13, gap=10)

    # Bottom: the ask
    box(slide, Inches(0.3), Inches(5.7), Inches(12.7), Inches(1.5), fill=C_DARK)
    txt(slide, "The ask",
        Inches(0.5), Inches(5.8), Inches(2.0), Inches(0.4),
        size=16, bold=True, color=C_AMBER)
    bullets(slide, [
        "Help inventorying existing bulk RNA-seq: which donors, tissues, and FASTQs are available?",
        "Pipeline is ready to scale \u2014 the bottleneck is organizing the data we already have",
    ], Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.8),
       size=13, gap=6, color=RGBColor(0xEE, 0xEE, 0xEE))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    prs = new_prs()
    build(prs)
    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")
