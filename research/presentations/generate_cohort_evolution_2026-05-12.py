"""
generate_cohort_evolution_2026-05-12.py
Slide deck: DRG Virome Pipeline — Cohort Evolution (all runs, start to finish)
2026-05-12 — Matthew Wild, TJP Lab, UT Dallas

Run from repo root:
    python3 research/presentations/generate_cohort_evolution_2026-05-12.py

Output:
    research/presentations/virome_cohort_evolution_2026-05-12.pptx
"""

import os
import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paths ──────────────────────────────────────────────────────────────────────
REPO_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
OUT_PATH  = os.path.join(REPO_ROOT, "research", "presentations",
                         "virome_cohort_evolution_2026-05-12.pptx")
HERVK_FIG = os.path.join(REPO_ROOT, "results", "hervk_cohort_comparison.png")

# ── Design constants ───────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

C_BG     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
C_BLUE   = RGBColor(0x2E, 0x86, 0xAB)
C_CORAL  = RGBColor(0xE8, 0x48, 0x55)
C_GREEN  = RGBColor(0x3B, 0xB2, 0x73)
C_AMBER  = RGBColor(0xF9, 0xC7, 0x4F)
C_PURPLE = RGBColor(0x7B, 0x2D, 0x8B)
C_TEAL   = RGBColor(0x00, 0x96, 0x88)
C_ORANGE = RGBColor(0xF5, 0x7C, 0x00)
C_LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
C_MGRAY  = RGBColor(0x9E, 0x9E, 0x9E)
C_STRIPE = RGBColor(0xE8, 0xF4, 0xFD)

TITLE_FONT = "Calibri"
BODY_FONT  = "Calibri"

# Cohort accent colors (one per cohort group)
COHORT_COLORS = {
    1:  C_BLUE,
    2:  C_TEAL,
    5:  C_GREEN,
    6:  C_CORAL,
    7:  C_ORANGE,
    8:  C_PURPLE,
    9:  C_AMBER,
    10: RGBColor(0x37, 0x47, 0x4F),
}

# ── Cohort data ────────────────────────────────────────────────────────────────
COHORTS = [
    {
        "num": 1, "name": "Muscle",
        "samples": "Sample_19–23", "n": 5,
        "tissue": "Skeletal muscle", "library": "PolyA, 150 bp PE",
        "pipeline": "v1.0", "sequencer": "Illumina",
        "virus": "None detected",
        "findings": [
            "Pipeline baseline run — first cohort processed end-to-end",
            "Established artifact exclusion list (k-mer cross-mapping from ruminant/insect viruses)",
            "HERV-K detected at low level (mean 31.3 RPM); confirmed pan-sample signal",
            "No human pathogens; clean tissue as expected for skeletal muscle",
            "Saad_1 equivalent: no QC outliers in this cohort",
        ],
    },
    {
        "num": 2, "name": "Early DRG (Donor1 · AIG1390 · Saad)",
        "samples": "donor1_L1–L5/T12, AIG1390_L1–L4/T12, Saad_1–5",
        "n": 16,
        "tissue": "DRG (post-mortem)", "library": "PolyA, 150 bp PE",
        "pipeline": "v1.0–1.2", "sequencer": "Illumina",
        "virus": "HERV-K (all), HHV-5 proxy (subset)",
        "findings": [
            "Three donors, multi-level sampling (L1–L5/T12) — first DRG virome profiles",
            "HERV-K confirmed as pan-sample DRG signal; forms baseline for all future cohorts",
            "HHV-5 proxy (Cytomegalovirus papiinebeta3 → remapped to 'Human CMV [HHV-5] proxy')",
            "Saad_1: QC outlier — 10× sequencing depth, extreme contamination profile; retained for pipeline assessment",
            "Saad_2: known library failure — retained for QC characterization",
            "Taxon display-name remapping (taxon_remap.tsv) introduced to address misleading NCBI labels",
        ],
    },
    {
        "num": 5, "name": "REJOIN Jayden",
        "samples": "473-1 – 473-17", "n": 17,
        "tissue": "DRG (post-mortem)", "library": "PolyA, 150 bp PE",
        "pipeline": "v1.2–1.3", "sequencer": "Illumina",
        "virus": "HERV-K (all), HHV-5 proxy (subset)",
        "findings": [
            "17-sample REJOIN cohort; completes the 38-sample full cohort used for Paper 1",
            "Paper 1 analysis: all_cohort_pluspf run consolidates cohorts 1–5",
            "Dual-database framework (v1.3) introduced: viral-only DB + PlusPF DB in parallel",
            "Three-tier output: Tier 1 (both DBs agree), Tier 2 (viral-only only), Tier 3 (PlusPF only)",
            "No novel human pathogens; HERV-K and CMV proxy confirmed across cohort",
        ],
    },
    {
        "num": 6, "name": "Parkinson 2026",
        "samples": "PD2–PD20, 023–028", "n": 20,
        "tissue": "DRG (post-mortem)", "library": "PolyA, 150 bp PE",
        "pipeline": "v1.5.0", "sequencer": "Illumina (Psomagen AN00028264)",
        "virus": "HSV-1 Tier 1 (PD19), HERV-K (all), HHV-5 proxy (subset)",
        "findings": [
            "14 Parkinson's disease patients + 6 non-PD controls; sequenced by Psomagen",
            "PD19: first-ever HSV-1 Tier 1 detection in the dataset — landmark finding",
            "PD1/7/8/11/12/13 absent from sequencing delivery (not run)",
            "No significant HERV-K difference between PD cases and controls (p=0.397)",
            "Dual-database validation: HSV-1 in PD19 confirmed Tier 1 (both DBs agree)",
            "Triggered blast_verify.nf offshoot to confirm identity and latency phase",
        ],
    },
    {
        "num": 7, "name": "BLAST Verify PD19",
        "samples": "PD19", "n": 1,
        "tissue": "DRG (post-mortem)", "library": "PolyA, 150 bp PE",
        "pipeline": "v1.5.0 (blast_verify.nf)", "sequencer": "Illumina",
        "virus": "HSV-1 confirmed",
        "findings": [
            "Offshoot pipeline (blast_verify.nf) — post-hoc re-analysis of PD19",
            "Extracted Kraken2-assigned HSV-1 reads → BLAST against NCBI nt database",
            "Confirmed identity: Human herpesvirus 1 (NC_001806.2), >98% identity",
            "Read distribution consistent with latency: enriched at LAT locus",
            "Established blast_verify.nf as reusable module for future Tier 1 candidates",
            "save_kraken2_output and save_unmapped_reads params added to main pipeline",
        ],
    },
    {
        "num": 8, "name": "Iadorola TG",
        "samples": "TG1–TG22 (non-contiguous)", "n": 16,
        "tissue": "Trigeminal ganglion (post-mortem)", "library": "Total RNA, 125 bp PE",
        "pipeline": "v1.5.0", "sequencer": "Illumina HiSeq 2500",
        "virus": "HSV-1 Tier 1 (5/16), HERV-K (all), HHV-5 proxy (subset)",
        "findings": [
            "Public dataset (LaPaglia et al. 2017, SRP113004) — pipeline validation benchmark",
            "HSV-1 detected in 5/16 samples: exact match to published results",
            "Two high-level samples (TG12: 14.7 RPM, TG3: 12.5 RPM) match paper's 'two high-level' report",
            "6th seropositive sample missed by both our pipeline and MAGIC — consistent non-detection",
            "Linear regression: HERV-K ~ HSV-1 RPM, R²=0.459, p=0.0039 (novel finding)",
            "VZV absent across all samples — consistent with known low LAT expression in TG",
        ],
    },
    {
        "num": 9, "name": "OSM Juliet",
        "samples": "D1–D3 × O/V × rep1–3", "n": 18,
        "tissue": "DRG (cultured, in vitro)", "library": "PolyA, 150 bp PE",
        "pipeline": "v1.5.0", "sequencer": "Illumina",
        "virus": "HERV-K only",
        "findings": [
            "3 donors × OSM vs vehicle × 3 replicates — first in vitro DRG cohort",
            "Clean virome: no human pathogens detected in any sample",
            "HERV-K significantly elevated vs all post-mortem cohorts (p<0.001)",
            "OSM treatment has no effect on HERV-K (OSM vs vehicle p=0.791)",
            "Culture context itself drives HERV-K elevation — stress/passaging, not cytokine",
            "Confirms Gihfavirus/Kinglevirus as post-mortem DRG artifact (absent in culture)",
        ],
    },
    {
        "num": 10, "name": "DPN & RA Kulkarni",
        "samples": "JE81, WE52, PA71, J1M2, T0M2, 2AC1, 1AN1, 8EN1, RO52, S1D1,\n2OE1, N4N2, ME61, LE41, L1V1, FL0A2, FL0B1, 1VYA1, 1VYB2,\nCA7A2, CA7B1, 4LIA1, 4LIB2, G1NA1, G1NB2",
        "n": 25,
        "tissue": "DRG (post-mortem)", "library": "PolyA, 150 bp PE",
        "pipeline": "v1.5.0", "sequencer": "Illumina",
        "virus": "TBD (pipeline running)",
        "findings": [
            "10 shared healthy controls + 5 DPN cases + 10 RA specimens (5 donors × A/B)",
            "First disease-comparative DRG cohort: diabetic peripheral neuropathy vs rheumatoid arthritis",
            "A/B specimen designation unconfirmed (L/R DRG or technical replicates — to clarify)",
            "5 metadata donors without FASTQ files (M4X, L3O, J4N, M3L, T1F — separate batch)",
            "Pipeline currently running on Juno — results pending",
            "Will enable first cross-disease virome comparison in human sensory ganglia",
        ],
    },
]

# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, x, y, w, h, fill=None, line=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.background() if fill is None else (shape.fill.solid(), setattr(shape.fill.fore_color, 'rgb', fill))
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line
        shape.line.width = line_width
    return shape

def tb(slide, text, x, y, w, h,
       font=BODY_FONT, size=Pt(13), bold=False, italic=False,
       color=C_DARK, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.name  = font
    r.font.size  = size
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def tb_multi(slide, lines, x, y, w, h,
             font=BODY_FONT, size=Pt(13), color=C_DARK,
             bold_first=False, line_spacing=1.25):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100 * 100)))
        r = p.add_run()
        r.text = line
        r.font.name  = font
        r.font.size  = size
        r.font.color.rgb = color
        r.font.bold  = (bold_first and i == 0)
    return box

def embed_png_bytes(slide, png_bytes, x, y, w, h):
    img_stream = io.BytesIO(png_bytes)
    slide.shapes.add_picture(img_stream, x, y, w, h)

def fig_to_bytes(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    buf.seek(0)
    return buf.read()

# ── Overview table figure ──────────────────────────────────────────────────────

def make_overview_table():
    rows = [
        ["#", "Cohort", "Tissue", "n", "Library", "Human Virus Detected", "Notable Finding"],
        ["1", "Muscle", "Skeletal muscle", "5", "PolyA 150PE",
         "None", "Pipeline baseline; artifact list seeded"],
        ["2–4", "Early DRG\n(Donor1/AIG1390/Saad)", "DRG (post-mortem)", "16", "PolyA 150PE",
         "HERV-K, HHV-5 proxy", "HERV-K baseline; Saad_1 QC outlier"],
        ["5", "REJOIN Jayden", "DRG (post-mortem)", "17", "PolyA 150PE",
         "HERV-K, HHV-5 proxy", "Completes 38-sample Paper 1 cohort"],
        ["6", "Parkinson 2026", "DRG (post-mortem)", "20", "PolyA 150PE",
         "HSV-1★, HERV-K, HHV-5 proxy", "First HSV-1 Tier 1 detection (PD19)"],
        ["7", "BLAST Verify PD19", "DRG (post-mortem)", "1", "PolyA 150PE",
         "HSV-1 confirmed", "LAT-phase confirmed; blast_verify.nf"],
        ["8", "Iadorola TG", "TG (post-mortem)", "16", "Total RNA 125PE",
         "HSV-1 (5/16)★, HERV-K", "Matches LaPaglia 2017; HERV-K~HSV-1 R²=0.46"],
        ["9", "OSM Juliet", "DRG (cultured)", "18", "PolyA 150PE",
         "HERV-K only", "HERV-K↑ in culture; no pathogens"],
        ["10", "DPN & RA Kulkarni", "DRG (post-mortem)", "25", "PolyA 150PE",
         "TBD (running)", "First DPN vs RA disease comparison"],
    ]

    col_widths = [0.32, 1.55, 1.35, 0.28, 0.95, 1.80, 2.55]
    total_w = sum(col_widths)
    fig_w, fig_h = total_w + 0.2, len(rows) * 0.38 + 0.15
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    fig.patch.set_facecolor("white")
    ax.set_xlim(0, total_w)
    ax.set_ylim(0, len(rows))
    ax.axis("off")

    header_color  = "#1A1A2E"
    stripe_color  = "#EBF5FB"
    finding_color = "#FEF9E7"
    text_dark     = "#1A1A2E"
    text_light    = "white"

    accent_map = {
        "1":   "#2E86AB",
        "2–4": "#009688",
        "5":   "#3BB273",
        "6":   "#E84855",
        "7":   "#F57C00",
        "8":   "#7B2D8B",
        "9":   "#F9C74F",
        "10":  "#37474F",
    }

    for row_idx, row in enumerate(rows):
        y_pos = len(rows) - row_idx - 1
        x_pos = 0
        is_header = (row_idx == 0)
        bg = header_color if is_header else (stripe_color if row_idx % 2 == 0 else "white")

        for col_idx, (cell, cw) in enumerate(zip(row, col_widths)):
            cell_bg = bg
            if not is_header and col_idx == 6:
                cell_bg = finding_color
            if not is_header and col_idx == 0:
                acc = accent_map.get(row[0], "#2E86AB")
                rect_patch = plt.Rectangle((x_pos, y_pos), cw, 1,
                                           facecolor=acc, edgecolor="white", linewidth=0.5)
                ax.add_patch(rect_patch)
                fc = "white"
            else:
                rect_patch = plt.Rectangle((x_pos, y_pos), cw, 1,
                                           facecolor=cell_bg, edgecolor="white", linewidth=0.5)
                ax.add_patch(rect_patch)
                fc = text_light if is_header else text_dark

            fs = 7.5 if not is_header else 8
            fw = "bold" if is_header or col_idx == 0 else "normal"
            ax.text(x_pos + cw / 2, y_pos + 0.5, cell,
                    ha="center", va="center", fontsize=fs,
                    fontweight=fw, color=fc, wrap=True,
                    multialignment="center",
                    linespacing=1.3)
            x_pos += cw

    # star legend
    ax.text(total_w - 0.01, 0.08, "★ = Tier 1 (dual-DB confirmed)",
            ha="right", va="bottom", fontsize=6.5, color="#555", style="italic")

    fig.tight_layout(pad=0.1)
    return fig_to_bytes(fig), fig_w, fig_h


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_title(prs):
    s = blank_slide(prs)
    rect(s, Inches(0), Inches(0), W, H, fill=C_DARK)
    rect(s, Inches(0), Inches(5.6), W, Inches(1.9), fill=C_BLUE)

    tb(s, "DRG Virome Pipeline",
       Inches(1), Inches(1.2), Inches(11.3), Inches(1.2),
       font=TITLE_FONT, size=Pt(48), bold=True, color=RGBColor(0xFF,0xFF,0xFF),
       align=PP_ALIGN.CENTER)
    tb(s, "Cohort Evolution: Start to Finish",
       Inches(1), Inches(2.55), Inches(11.3), Inches(0.8),
       font=TITLE_FONT, size=Pt(30), bold=False, color=C_AMBER,
       align=PP_ALIGN.CENTER)

    tb(s, "TJP Group · UT Dallas · 2026",
       Inches(1), Inches(5.75), Inches(6), Inches(0.6),
       font=BODY_FONT, size=Pt(16), color=RGBColor(0xFF,0xFF,0xFF))
    tb(s, "DRG Virome Pipeline v1.5.0\n10 cohorts · 112 ganglion samples · 2 tissue types",
       Inches(1), Inches(6.3), Inches(11), Inches(0.8),
       font=BODY_FONT, size=Pt(14), color=RGBColor(0xFF,0xFF,0xFF))
    return s


def slide_overview(prs):
    s = blank_slide(prs)
    rect(s, Inches(0), Inches(0), W, Inches(0.85), fill=C_DARK)
    tb(s, "Cohort Overview",
       Inches(0.3), Inches(0.12), Inches(9), Inches(0.6),
       font=TITLE_FONT, size=Pt(24), bold=True,
       color=RGBColor(0xFF,0xFF,0xFF))
    tb(s, "All pipeline runs · final data only · non-contiguous cohort numbers reflect registry order",
       Inches(0.3), Inches(0.55), Inches(12), Inches(0.35),
       font=BODY_FONT, size=Pt(11), color=C_AMBER)

    png_bytes, fw, fh = make_overview_table()
    scale = min(Inches(12.7) / Inches(fw), Inches(6.3) / Inches(fh))
    iw = Inches(fw) * scale / Inches(1)
    ih = Inches(fh) * scale / Inches(1)
    ix = (Inches(13.33) - Inches(iw)) / 2
    iy = Inches(0.95)
    embed_png_bytes(s, png_bytes, Inches(ix / Inches(1)), iy, Inches(iw), Inches(ih))
    return s


def slide_cohort(prs, cohort):
    s    = blank_slide(prs)
    acc  = COHORT_COLORS.get(cohort["num"], C_BLUE)

    # Header bar
    rect(s, Inches(0), Inches(0), W, Inches(1.0), fill=C_DARK)
    rect(s, Inches(0), Inches(0), Inches(0.45), Inches(1.0), fill=acc)

    num_str = str(cohort["num"])
    tb(s, num_str,
       Inches(0), Inches(0.1), Inches(0.45), Inches(0.8),
       font=TITLE_FONT, size=Pt(28), bold=True,
       color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)

    tb(s, cohort["name"],
       Inches(0.55), Inches(0.15), Inches(9), Inches(0.65),
       font=TITLE_FONT, size=Pt(26), bold=True,
       color=RGBColor(0xFF,0xFF,0xFF))
    tb(s, f"{cohort['tissue']}  ·  n={cohort['n']}  ·  {cohort['pipeline']}",
       Inches(0.55), Inches(0.68), Inches(9), Inches(0.3),
       font=BODY_FONT, size=Pt(11), color=C_AMBER)

    # Left panel — metadata
    rect(s, Inches(0), Inches(1.0), Inches(4.1), Inches(6.5), fill=C_LGRAY)
    rect(s, Inches(0), Inches(1.0), Inches(0.06), Inches(6.5), fill=acc)

    meta_labels = ["Samples", "n", "Tissue", "Library", "Sequencer", "Pipeline", "Virus Detected"]
    meta_values = [
        cohort["samples"], str(cohort["n"]), cohort["tissue"],
        cohort["library"], cohort["sequencer"], cohort["pipeline"],
        cohort["virus"],
    ]

    y0 = Inches(1.18)
    row_h = Inches(0.68)
    for i, (lbl, val) in enumerate(zip(meta_labels, meta_values)):
        y = y0 + i * row_h
        tb(s, lbl,
           Inches(0.18), y, Inches(1.1), row_h,
           size=Pt(10), bold=True, color=C_MGRAY)
        tb(s, val,
           Inches(0.18), y + Inches(0.22), Inches(3.8), row_h,
           size=Pt(11), bold=False, color=C_DARK)

    # Right panel — findings
    tb(s, "Key Findings",
       Inches(4.3), Inches(1.05), Inches(8.8), Inches(0.45),
       font=TITLE_FONT, size=Pt(16), bold=True, color=acc)

    bullet_lines = [f"  •  {f}" for f in cohort["findings"]]
    tb_multi(s, bullet_lines,
             Inches(4.3), Inches(1.55), Inches(8.8), Inches(5.7),
             size=Pt(13.5), color=C_DARK, line_spacing=1.4)

    # Bottom bar
    rect(s, Inches(0), Inches(7.28), W, Inches(0.22), fill=acc)
    tb(s, f"DRG Virome Pipeline  ·  Cohort {cohort['num']} of 10  ·  TJP Group UT Dallas",
       Inches(0.2), Inches(7.28), Inches(12), Inches(0.22),
       size=Pt(8), color=RGBColor(0xFF,0xFF,0xFF))
    return s


def slide_hervk(prs):
    """Extra slide: HERV-K cross-cohort comparison figure."""
    if not os.path.exists(HERVK_FIG):
        return
    s = blank_slide(prs)
    rect(s, Inches(0), Inches(0), W, Inches(1.0), fill=C_DARK)
    rect(s, Inches(0), Inches(0), Inches(0.45), Inches(1.0),
         fill=RGBColor(0x43, 0xA0, 0x47))
    tb(s, "★",
       Inches(0), Inches(0.1), Inches(0.45), Inches(0.8),
       font=TITLE_FONT, size=Pt(28), bold=True,
       color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    tb(s, "Cross-Cohort Analysis: HERV-K Expression",
       Inches(0.55), Inches(0.15), Inches(10), Inches(0.65),
       font=TITLE_FONT, size=Pt(26), bold=True,
       color=RGBColor(0xFF,0xFF,0xFF))
    tb(s, "Kruskal-Wallis H=31.71, p=0.000002  ·  Pairwise Mann-Whitney with Bonferroni correction",
       Inches(0.55), Inches(0.68), Inches(11), Inches(0.3),
       font=BODY_FONT, size=Pt(11), color=C_AMBER)

    s.shapes.add_picture(HERVK_FIG, Inches(0.3), Inches(1.1), Inches(8.5), Inches(5.8))

    findings = [
        "OSM Juliet (cultured DRG) significantly\nhigher than all post-mortem cohorts",
        "Effect is culture context, not OSM\ntreatment (OSM vs vehicle: p=0.79)",
        "Post-mortem tissue consistent: Muscle,\nPD DRG, TG statistically equivalent",
        "No HERV-K elevation in PD vs controls\n(p=0.40) at current n",
        "HERV-K > 50 RPM → suspect in vitro\nprocessing (QC anchor for new cohorts)",
    ]
    y0 = Inches(1.15)
    for i, f in enumerate(findings):
        rect(s, Inches(9.0), y0 + i*Inches(1.08), Inches(4.1), Inches(0.95),
             fill=C_STRIPE, line=C_BLUE, line_width=Pt(0.75))
        tb(s, f, Inches(9.12), y0 + i*Inches(1.08) + Inches(0.08),
           Inches(3.9), Inches(0.9), size=Pt(11), color=C_DARK)

    rect(s, Inches(0), Inches(7.28), W, Inches(0.22),
         fill=RGBColor(0x43, 0xA0, 0x47))
    tb(s, "DRG Virome Pipeline  ·  Cross-cohort HERV-K analysis  ·  TJP Group UT Dallas",
       Inches(0.2), Inches(7.28), Inches(12), Inches(0.22),
       size=Pt(8), color=RGBColor(0xFF,0xFF,0xFF))


# ── Build deck ─────────────────────────────────────────────────────────────────

prs = new_prs()

slide_title(prs)
slide_overview(prs)
for c in COHORTS:
    slide_cohort(prs, c)
slide_hervk(prs)

prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
print(f"Slides: {len(prs.slides)}")
