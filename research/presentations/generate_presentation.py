"""
generate_presentation.py
Lab presentation: Systematic Virome Profiling of Human Neural Tissue
2026-03-28 — Matthew Wild, TJP Lab, UT Dallas

Run from repo root:
    python3 research/presentations/generate_presentation.py

Outputs:
    research/presentations/virome_lab_presentation_2026-03-28.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
REPO_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
FIG_DIR = os.path.join(REPO_ROOT, "research", "paper1", "figures")
OUT_PATH = os.path.join(REPO_ROOT, "research", "presentations", "virome_lab_presentation_2026-03-28.pptx")

FIG1 = os.path.join(FIG_DIR, "fig1_pipeline_diagram.png")
FIG2 = os.path.join(FIG_DIR, "fig2_filtering_funnel.png")
FIG3 = os.path.join(FIG_DIR, "fig3_tier_summary.png")
FIG4 = os.path.join(FIG_DIR, "fig4_herv_k.png")

# ---------------------------------------------------------------------------
# Design constants (16:9 = 13.33" × 7.5")
# ---------------------------------------------------------------------------
W = Inches(13.33)
H = Inches(7.5)

# Color palette
C_BG       = RGBColor(0xFF, 0xFF, 0xFF)   # white
C_DARK     = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
C_BLUE     = RGBColor(0x2E, 0x86, 0xAB)   # steel blue — method
C_CORAL    = RGBColor(0xE8, 0x48, 0x55)   # coral/red — key findings
C_GREEN    = RGBColor(0x3B, 0xB2, 0x73)   # green — positive/confirmed
C_AMBER    = RGBColor(0xF9, 0xC7, 0x4F)   # amber — warning/FP
C_PURPLE   = RGBColor(0x7B, 0x2D, 0x8B)   # purple — biological insight
C_LGRAY    = RGBColor(0xF5, 0xF5, 0xF5)   # light gray panel
C_MGRAY    = RGBColor(0x9E, 0x9E, 0x9E)   # medium gray
C_STRIPE   = RGBColor(0xE8, 0xF4, 0xFD)   # light blue stripe

TITLE_FONT   = "Calibri"
BODY_FONT    = "Calibri"
CODE_FONT    = "Courier New"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def rect(slide, x, y, w, h, fill=None, line=None, line_width=Pt(1)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.fill.background()
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line
        shape.line.width = line_width
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_name=BODY_FONT, font_size=Pt(14), bold=False,
                 italic=False, color=C_DARK, align=PP_ALIGN.LEFT,
                 word_wrap=True, line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if line_spacing:
        from pptx.util import Pt as Pt2
        from pptx.oxml.ns import qn as qn2
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn2('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn2('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txBox


def add_multiline(slide, lines, x, y, w, h,
                  font_name=BODY_FONT, font_size=Pt(13), bold=False,
                  color=C_DARK, line_gap_pt=6, align=PP_ALIGN.LEFT):
    """Add a text box with multiple paragraphs (list of strings or (text, bold, color) tuples)."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, b, col = item, bold, color
        else:
            txt, b, col = item[0], item[1] if len(item) > 1 else bold, item[2] if len(item) > 2 else color
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        # gap
        from pptx.oxml.ns import qn as qn2
        pPr = p._p.get_or_add_pPr()
        spcBef = etree.SubElement(pPr, qn2('a:spcBef'))
        spcPts = etree.SubElement(spcBef, qn2('a:spcPts'))
        spcPts.set('val', str(int(line_gap_pt * 100)))
        run = p.add_run()
        run.text = txt
        run.font.name = font_name
        run.font.size = font_size
        run.font.bold = b
        run.font.color.rgb = col
    return txBox


def slide_header(slide, title, subtitle=None,
                 bar_color=C_BLUE, title_color=C_DARK):
    """Adds a top bar + title."""
    # color bar
    rect(slide, 0, 0, W, Inches(0.08), fill=bar_color)
    # title
    add_text_box(slide, title,
                 x=Inches(0.45), y=Inches(0.12), w=Inches(12.0), h=Inches(0.7),
                 font_size=Pt(28), bold=True, color=title_color)
    if subtitle:
        add_text_box(slide, subtitle,
                     x=Inches(0.45), y=Inches(0.78), w=Inches(12.0), h=Inches(0.35),
                     font_size=Pt(14), color=C_MGRAY, italic=True)
    # bottom rule
    rect(slide, 0, Inches(1.08), W, Inches(0.025), fill=C_LGRAY)


def section_divider(prs, number, title, subtitle, color=C_BLUE):
    """Full-bleed section separator slide."""
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, fill=color)
    # section number pill
    pill = slide.shapes.add_shape(1, Inches(0.5), Inches(2.8), Inches(1.1), Inches(0.55))
    pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
    pill.line.fill.background()
    add_text_box(slide, f"0{number}", Inches(0.5), Inches(2.78), Inches(1.1), Inches(0.6),
                 font_size=Pt(22), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text_box(slide, title,
                 Inches(0.5), Inches(3.45), Inches(12.0), Inches(1.2),
                 font_size=Pt(46), bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    add_text_box(slide, subtitle,
                 Inches(0.5), Inches(4.65), Inches(12.0), Inches(0.7),
                 font_size=Pt(18), color=RGBColor(0xFF,0xFF,0xFF), italic=True)
    return slide


def add_figure(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        # placeholder box
        r = rect(slide, x, y, w, h or Inches(3), fill=C_LGRAY, line=C_MGRAY)
        add_text_box(slide, f"[Figure: {os.path.basename(path)}]",
                     x + Inches(0.1), y + Inches(0.1), w - Inches(0.2), Inches(0.5),
                     font_size=Pt(11), color=C_MGRAY)
        return
    if h:
        slide.shapes.add_picture(path, x, y, width=w, height=h)
    else:
        slide.shapes.add_picture(path, x, y, width=w)


def bullet_box(slide, items, x, y, w, h, font_size=Pt(13.5), gap=8, icon="▸", color=None):
    """
    items: list of strings or (text, indent_level) tuples.
    indent_level 0 = main bullet, 1 = sub-bullet
    """
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            txt, level = item, 0
        else:
            txt, level = item[0], item[1]
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        from pptx.oxml.ns import qn as qn2
        pPr = p._p.get_or_add_pPr()
        spcBef = etree.SubElement(pPr, qn2('a:spcBef'))
        spcPts = etree.SubElement(spcBef, qn2('a:spcPts'))
        spcPts.set('val', str(gap * 100))
        indent_str = "  " * level
        run = p.add_run()
        prefix = "  ◦ " if level > 0 else f"{icon} "
        run.text = indent_str + prefix + txt
        run.font.name = BODY_FONT
        run.font.size = font_size
        run.font.color.rgb = color if color is not None else (C_DARK if level == 0 else C_MGRAY)
    return txBox


def stat_box(slide, value, label, x, y, w=Inches(2.0), h=Inches(1.2),
             val_color=C_BLUE, bg=C_LGRAY):
    rect(slide, x, y, w, h, fill=bg)
    add_text_box(slide, value, x, y + Inches(0.1), w, Inches(0.65),
                 font_size=Pt(30), bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, x, y + Inches(0.7), w, Inches(0.45),
                 font_size=Pt(11), color=C_DARK, align=PP_ALIGN.CENTER, word_wrap=True)


def footnote(slide, text, y=Inches(7.1)):
    add_text_box(slide, text, Inches(0.45), y, Inches(12.4), Inches(0.35),
                 font_size=Pt(9), color=C_MGRAY, italic=True)


# ---------------------------------------------------------------------------
# Build slides
# ---------------------------------------------------------------------------

def build(prs):

    # ── SLIDE 1: TITLE ──────────────────────────────────────────────────────
    slide = blank_slide(prs)
    # gradient-ish block
    rect(slide, 0, 0, W, Inches(5.2), fill=C_DARK)
    rect(slide, 0, Inches(5.2), W, Inches(2.3), fill=C_BLUE)

    add_text_box(slide,
        "Systematic Virome Profiling\nof Human Neural Tissue",
        Inches(0.6), Inches(0.7), Inches(10.5), Inches(2.4),
        font_size=Pt(42), bold=True,
        color=RGBColor(0xFF,0xFF,0xFF))
    add_text_box(slide,
        "A Nextflow pipeline for competitive dual-database classification from bulk RNA-seq",
        Inches(0.6), Inches(3.1), Inches(10.5), Inches(0.6),
        font_size=Pt(18), color=C_AMBER, italic=False)
    add_text_box(slide,
        "Matthew Wild  ·  TJP Lab  ·  UT Dallas  ·  March 2026",
        Inches(0.6), Inches(3.75), Inches(10.5), Inches(0.5),
        font_size=Pt(14), color=RGBColor(0xCC,0xCC,0xCC))

    add_text_box(slide,
        "What this talk covers:",
        Inches(0.6), Inches(5.35), Inches(5.5), Inches(0.4),
        font_size=Pt(13), bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    items = [
        "The method — solving sensitivity vs. specificity in virome profiling",
        "The findings — what our data actually contains (and doesn't)",
        "The standard — how this aligns with cutting-edge genomics practice",
        "The ask — consolidating the lab's sequencing assets at scale",
    ]
    for i, item in enumerate(items):
        add_text_box(slide,
            f"{'①②③④'[i]}  {item}",
            Inches(0.6), Inches(5.75 + i * 0.38), Inches(11.5), Inches(0.38),
            font_size=Pt(13), color=RGBColor(0xFF,0xFF,0xFF))


    # ── SLIDE 2: WHY VIRUSES IN DRG ─────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Why look for viruses in dorsal root ganglia?",
                 subtitle="Neurotropic viruses establish latency in sensory ganglia — and may drive pain long after acute infection",
                 bar_color=C_DARK)

    # Three columns
    col_w = Inches(3.9)
    col_h = Inches(4.8)
    cols = [
        {
            "title": "The Biology",
            "color": C_BLUE,
            "items": [
                "DRG = primary sensory neurons; peripheral pain pathway entry point",
                "VZV (varicella-zoster) establishes latency in DRG — reactivation → shingles / postherpetic neuralgia",
                "HSV-1/2 latency in trigeminal & sacral DRG",
                "CMV, HHV-6 detected in neural tissue under immunosuppression",
                "HIV — peripheral neuropathy in 30–50% of patients",
                "SARS-CoV-2 long COVID: sensory neuropathy, small fiber neuropathy",
            ]
        },
        {
            "title": "The Clinical Need",
            "color": C_CORAL,
            "items": [
                "Chronic neuropathic pain affects ~8% of adults — mechanisms poorly understood",
                "Viral triggers implicated in: PHN, DPN, HIV neuropathy, MS nerve damage",
                "No systematic profiling of the DRG virome at transcriptomic resolution",
                "Human DRG tissue increasingly available via post-mortem & surgical programs",
                "GTEx, NABEC, and disease-specific biobanks now provide bulk RNA-seq data",
            ]
        },
        {
            "title": "The Opportunity",
            "color": C_GREEN,
            "items": [
                "Our lab has 100–200 human donors across 12+ tissue types",
                "Bulk RNA-seq already in hand for DRG, spinal cord, muscle, and more",
                "Single-nucleus, spatial, and multi-omic modalities available for DRG",
                "Computational pipeline can answer this question at scale — if we use the data we have",
                "Paper 2 target: neuropathic pain vs. healthy DRG across all pain conditions",
            ]
        }
    ]
    for i, col in enumerate(cols):
        x = Inches(0.3 + i * 4.5)
        rect(slide, x, Inches(1.25), col_w, Inches(0.45), fill=col["color"])
        add_text_box(slide, col["title"],
                     x + Inches(0.1), Inches(1.27), col_w - Inches(0.2), Inches(0.4),
                     font_size=Pt(14), bold=True,
                     color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
        rect(slide, x, Inches(1.7), col_w, col_h, fill=C_LGRAY)
        bullet_box(slide, col["items"],
                   x + Inches(0.15), Inches(1.78), col_w - Inches(0.3), col_h - Inches(0.15),
                   font_size=Pt(12), gap=5)


    # ── SLIDE 3: THE NEEDLE-IN-A-HAYSTACK PROBLEM ────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "The challenge: finding viral signal in a sea of host RNA",
                 subtitle="Bulk RNA-seq is not designed for metagenomics — but it's what we have",
                 bar_color=C_CORAL)

    # Left: composition visualization
    rect(slide, Inches(0.35), Inches(1.3), Inches(5.8), Inches(5.8), fill=C_LGRAY)

    add_text_box(slide, "Typical bulk RNA-seq library composition",
                 Inches(0.5), Inches(1.35), Inches(5.5), Inches(0.4),
                 font_size=Pt(13), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    # Stacked bar visual using shapes
    bar_x = Inches(1.2)
    bar_y = Inches(1.9)
    bar_w = Inches(3.5)

    # Host: 99.95%
    host_h = Inches(4.7)
    rect(slide, bar_x, bar_y, bar_w, host_h, fill=RGBColor(0xBB,0xDE,0xFB))
    add_text_box(slide, "Host (human)\n~99.95% of reads",
                 bar_x + Inches(0.1), bar_y + Inches(2.0), bar_w - Inches(0.2), Inches(0.7),
                 font_size=Pt(13), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    # Viral: 0.05% (visible sliver at top)
    rect(slide, bar_x, bar_y, bar_w, Inches(0.22), fill=C_CORAL)
    add_text_box(slide, "Viral reads: <0.05%",
                 bar_x + Inches(0.1), bar_y + Inches(0.23), bar_w, Inches(0.35),
                 font_size=Pt(11), color=C_CORAL, bold=True)

    # Right: the two failure modes
    rect(slide, Inches(6.4), Inches(1.3), Inches(6.55), Inches(5.8), fill=C_LGRAY)
    add_text_box(slide, "Two failure modes for virome classification",
                 Inches(6.55), Inches(1.35), Inches(6.2), Inches(0.4),
                 font_size=Pt(13), bold=True, color=C_DARK)

    # Failure mode 1
    rect(slide, Inches(6.55), Inches(1.85), Inches(6.1), Inches(0.4), fill=C_CORAL)
    add_text_box(slide, "Problem 1: Viral-only database — the closed-world assumption",
                 Inches(6.65), Inches(1.87), Inches(6.0), Inches(0.38),
                 font_size=Pt(12), bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    bullet_box(slide,
        ["Classifier has no non-viral reference — every read must land somewhere viral",
         "Human transcripts sharing k-mers with viral genomes → misassigned",
         "HERVs, CMV k-mers, and neural lncRNAs are common offenders",
         "Result: systematic false positives in every sample",
        ],
        Inches(6.55), Inches(2.3), Inches(6.1), Inches(1.85),
        font_size=Pt(12), gap=5)

    # Failure mode 2
    rect(slide, Inches(6.55), Inches(4.3), Inches(6.1), Inches(0.4), fill=C_AMBER)
    add_text_box(slide, "Problem 2: Sensitivity — latent viruses leave few traces",
                 Inches(6.65), Inches(4.32), Inches(6.0), Inches(0.38),
                 font_size=Pt(12), bold=True, color=C_DARK)
    bullet_box(slide,
        ["VZV/HSV maintain latency with minimal transcription — may fall below noise floor",
         "Bulk RNA-seq detection floor: ~10 reads per million host-depleted reads",
         "Targeted enrichment (VirCapSeq) or pseudobulk needed for latent detection",
        ],
        Inches(6.55), Inches(4.75), Inches(6.1), Inches(1.5),
        font_size=Pt(12), gap=5)

    footnote(slide, "Detection floor ~10 RPM established empirically from this cohort")


    # ── SECTION DIVIDER 1 ───────────────────────────────────────────────────
    section_divider(prs, 1,
        "The Method",
        "virome-pipeline: competitive dual-database classification for human neural tissue",
        color=C_BLUE)


    # ── SLIDE 4: PIPELINE ARCHITECTURE ──────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "virome-pipeline v1.3.0 — architecture",
                 subtitle="Seven-step Nextflow DSL2 workflow with optional dual-database competitive classification",
                 bar_color=C_BLUE)
    add_figure(slide, FIG1,
               x=Inches(0.35), y=Inches(1.25), w=Inches(12.6))
    footnote(slide, "STAR host depletion → Kraken2 classification → Bracken re-estimation → multi-stage filtering → cohort aggregation  |  All steps containerized (Apptainer); SLURM-ready")


    # ── SLIDE 5: THE DUAL-DATABASE INSIGHT ──────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "The core innovation: competitive dual-database classification",
                 subtitle="Ask two questions simultaneously — does this virus survive competition with the full genome?",
                 bar_color=C_BLUE)

    # Left panel: DB1 alone
    rect(slide, Inches(0.35), Inches(1.25), Inches(3.8), Inches(5.75), fill=C_LGRAY)
    rect(slide, Inches(0.35), Inches(1.25), Inches(3.8), Inches(0.5), fill=C_CORAL)
    add_text_box(slide, "Database 1: Viral-only",
                 Inches(0.45), Inches(1.27), Inches(3.6), Inches(0.45),
                 font_size=Pt(14), bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    bullet_box(slide,
        ["~10,000 viral genomes",
         "No non-viral sequences",
         "Closed-world: every read forced into a viral bin",
         "Result: detects HERV-K, CMV proxy, MCV as 'viruses'",
         "False positive rate: 100% in this cohort",
        ],
        Inches(0.45), Inches(1.85), Inches(3.6), Inches(3.5),
        font_size=Pt(12.5), gap=6)

    # Middle: three tiers
    rect(slide, Inches(4.4), Inches(1.25), Inches(4.5), Inches(5.75), fill=C_DARK)
    add_text_box(slide, "Three-tier confidence scoring",
                 Inches(4.5), Inches(1.32), Inches(4.3), Inches(0.45),
                 font_size=Pt(14), bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)

    tiers = [
        (C_GREEN,  "Tier 1 — Shared",    "Detected in BOTH databases\nHigh-confidence viral candidate\nUse for biology",      "= 0 in our data"),
        (C_AMBER,  "Tier 2 — Viral-only","Viral-only DB exclusive\nFalse positive candidates\nInvestigate before interpreting", "= 3 taxa (all resolved as FP)"),
        (C_MGRAY,  "Tier 3 — PlusPF only","PlusPF-only\nNon-viral contaminants\nQC utility — not biology",               "= 10–179 taxa / sample"),
    ]
    for i, (col, title, desc, result) in enumerate(tiers):
        y = Inches(1.95 + i * 1.65)
        rect(slide, Inches(4.55), y, Inches(4.2), Inches(1.45), fill=col)
        add_text_box(slide, title,
                     Inches(4.65), y + Inches(0.08), Inches(2.5), Inches(0.38),
                     font_size=Pt(13), bold=True, color=C_DARK)
        add_text_box(slide, desc,
                     Inches(4.65), y + Inches(0.48), Inches(2.5), Inches(0.9),
                     font_size=Pt(10.5), color=C_DARK)
        add_text_box(slide, result,
                     Inches(7.0), y + Inches(0.08), Inches(1.6), Inches(0.9),
                     font_size=Pt(11), bold=True, color=C_DARK, align=PP_ALIGN.RIGHT)

    # Right panel: DB2
    rect(slide, Inches(9.15), Inches(1.25), Inches(3.8), Inches(5.75), fill=C_LGRAY)
    rect(slide, Inches(9.15), Inches(1.25), Inches(3.8), Inches(0.5), fill=C_GREEN)
    add_text_box(slide, "Database 2: PlusPF (standard)",
                 Inches(9.25), Inches(1.27), Inches(3.6), Inches(0.45),
                 font_size=Pt(14), bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    bullet_box(slide,
        ["Full bacterial, archaeal, viral, fungal + human reference",
         "~100 GB — human genome included",
         "Open-world: reads can land on the correct host sequence",
         "HERV-K → reclassified to Homo sapiens",
         "CMV proxy → 0 reads (k-mer artifact resolved)",
        ],
        Inches(9.25), Inches(1.85), Inches(3.6), Inches(3.5),
        font_size=Pt(12.5), gap=6)


    # ── SLIDE 6: FILTERING FUNNEL ────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Multi-stage filtering funnel — 15 samples, 3 tissue groups",
                 subtitle="Systematic noise reduction: Bracken raw → min-reads threshold → curated artifact exclusion",
                 bar_color=C_BLUE)
    add_figure(slide, FIG2,
               x=Inches(0.6), y=Inches(1.25), w=Inches(12.1))
    footnote(slide, "24-entry artifact exclusion list | categories: reagent contaminants, ruminant bunyaviruses, insect baculoviruses, phages, environmental metagenome viruses, DRG k-mer cross-mapping")


    # ── SLIDE 7: ARTIFACT EXCLUSION FRAMEWORK ────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Curated artifact exclusion — 24 entries, versioned with the pipeline",
                 subtitle="Every false positive was investigated at the read level before being added to the exclusion list",
                 bar_color=C_BLUE)

    # Table of artifact categories
    headers = ["Category", "Example taxa", "Mechanism", "Diagnostic signature"]
    rows = [
        ("Reagent contaminants",      "PhiX174, lambda phage",                   "Library prep spike-in k-mers",            "All samples, all tissues"),
        ("Ruminant orthobunyaviruses","Schmallenberg virus, Akabane orthobunyavirus","Cattle/sheep host k-mer overlap",       "DRG-enriched (neuronal transcripts)"),
        ("Insect baculoviruses",       "Autographa californica MNPV",             "Invertebrate cell line contaminants",     "Sporadic; reagent-associated"),
        ("Phages (ICTV-reclassified)", "Ralstonia phage p12J → Porrectionivirus", "Taxon ID superseded; old ID escapes list","Run after DB update = reappears"),
        ("Environmental metagenome",  "Gihfavirus, Kinglevirus",                  "DRG-specific lncRNA / ion channel k-mers","DRG-only; absent from muscle"),
        ("Giant amoeba viruses",       "Colossusvirus, Pandoravirus",             "Environmental metagenome contamination",  "Sporadic, sample-independent"),
        ("Hantaviruses (Oxbow)",       "Orthohantavirus oxbowense (2 taxon IDs)", "Confirmed k-mer cross-mapping",           "Present in ALL samples/tissues"),
    ]

    col_widths = [Inches(2.3), Inches(3.0), Inches(3.5), Inches(4.0)]
    header_y = Inches(1.3)
    x_starts = [Inches(0.3), Inches(2.65), Inches(5.7), Inches(9.25)]

    for j, (hdr, xw) in enumerate(zip(headers, col_widths)):
        rect(slide, x_starts[j], header_y, xw - Inches(0.05), Inches(0.42), fill=C_BLUE)
        add_text_box(slide, hdr, x_starts[j] + Inches(0.05), header_y + Inches(0.04),
                     xw - Inches(0.1), Inches(0.35),
                     font_size=Pt(11), bold=True, color=RGBColor(0xFF,0xFF,0xFF))

    for i, row in enumerate(rows):
        bg = C_LGRAY if i % 2 == 0 else RGBColor(0xFF,0xFF,0xFF)
        row_y = Inches(1.75 + i * 0.72)
        for j, (cell, xw) in enumerate(zip(row, col_widths)):
            rect(slide, x_starts[j], row_y, xw - Inches(0.05), Inches(0.68), fill=bg)
            add_text_box(slide, cell,
                         x_starts[j] + Inches(0.07), row_y + Inches(0.05),
                         xw - Inches(0.14), Inches(0.6),
                         font_size=Pt(10.5), color=C_DARK, word_wrap=True)

    footnote(slide, "ICTV reclassification caveat: after each database update, audit for taxon ID supersession. Both old and new IDs listed where applicable.")


    # ── SECTION DIVIDER 2 ───────────────────────────────────────────────────
    section_divider(prs, 2,
        "The Findings",
        "Zero confirmed viruses — and why that's the most important result",
        color=C_PURPLE)


    # ── SLIDE 8: COHORT OVERVIEW ─────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Study cohort — 15 human samples, two tissue types",
                 subtitle="5 skeletal muscle  ·  6 DRG (donor1, L1–L5 + T12)  ·  4 DRG (Saad cohort, independent donors)",
                 bar_color=C_PURPLE)

    # Three cohort boxes
    cohorts = [
        {
            "name": "Skeletal Muscle",
            "n": "n = 5",
            "color": C_BLUE,
            "samples": "Sample_19, 20, 21, 22, 23",
            "notes": ["Independent donors", "Post-mortem vastus lateralis", "Non-neural control tissue", "QC: all pass"],
        },
        {
            "name": "DRG — donor1",
            "n": "n = 6",
            "color": C_PURPLE,
            "samples": "L1, L2, L3, L4, L5, T12",
            "notes": ["Single donor; all spinal levels", "Intra-donor replication", "Allows spinal level stratification", "QC: all pass"],
        },
        {
            "name": "DRG — Saad cohort",
            "n": "n = 4",
            "color": C_GREEN,
            "samples": "Saad_1 (QC flag), Saad_3, Saad_4, Saad_5",
            "notes": ["Independent donors; published cohort", "Saad_1: 430K bracken_raw reads (8× DRG median) — flagged, included", "Saad_2: excluded (library failure, 0 final taxa)", "AIG1390: excluded (confirmed donor1 duplicate by MD5)"],
        },
    ]
    for i, c in enumerate(cohorts):
        x = Inches(0.3 + i * 4.35)
        rect(slide, x, Inches(1.3), Inches(4.15), Inches(5.8), fill=C_LGRAY)
        rect(slide, x, Inches(1.3), Inches(4.15), Inches(0.55), fill=c["color"])
        add_text_box(slide, c["name"],
                     x + Inches(0.1), Inches(1.32), Inches(2.6), Inches(0.48),
                     font_size=Pt(15), bold=True, color=RGBColor(0xFF,0xFF,0xFF))
        add_text_box(slide, c["n"],
                     x + Inches(2.8), Inches(1.32), Inches(1.25), Inches(0.48),
                     font_size=Pt(22), bold=True, color=RGBColor(0xFF,0xFF,0xFF),
                     align=PP_ALIGN.RIGHT)
        add_text_box(slide, c["samples"],
                     x + Inches(0.1), Inches(1.92), Inches(3.9), Inches(0.4),
                     font_size=Pt(11), color=C_DARK, italic=True)
        bullet_box(slide, c["notes"],
                   x + Inches(0.1), Inches(2.38), Inches(3.9), Inches(3.0),
                   font_size=Pt(12), gap=5)

    # Key exclusions callout
    rect(slide, Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.35), fill=RGBColor(0xFF,0xF0,0xCC))
    add_text_box(slide,
        "⚠  AIG1390 (5 samples) excluded after MD5 checksum confirmed identical to donor1 FASTQs — duplicate submission under two donor IDs. Correct AIG1390 FASTQs to be obtained from provider.",
        Inches(0.45), Inches(7.07), Inches(12.5), Inches(0.3),
        font_size=Pt(10), color=C_DARK, bold=False)


    # ── SLIDE 9: THE NULL RESULT ─────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "The result: zero confirmed viral detections — cohort-wide",
                 subtitle="Tier 1 = 0 across all 15 samples and both tissue types  ·  100% false positive rate for viral-only Kraken2",
                 bar_color=C_CORAL)

    # Big stat boxes
    stat_box(slide, "0",  "Tier 1 taxa\n(confirmed viral)", Inches(0.4),  Inches(1.3), Inches(2.3), Inches(1.5), val_color=C_GREEN)
    stat_box(slide, "3",  "Tier 2 taxa\n(all resolved as FP)", Inches(2.85), Inches(1.3), Inches(2.3), Inches(1.5), val_color=C_CORAL)
    stat_box(slide, "100%", "False positive rate\nof viral-only DB",  Inches(5.3),  Inches(1.3), Inches(2.3), Inches(1.5), val_color=C_CORAL)
    stat_box(slide, "15/15","Samples with\nTier 1 = 0",             Inches(7.75), Inches(1.3), Inches(2.3), Inches(1.5), val_color=C_GREEN)
    stat_box(slide, "~10 RPM","Detection floor\n(host-depleted fraction)", Inches(10.2), Inches(1.3), Inches(2.75), Inches(1.5), val_color=C_AMBER)

    # Narrative text
    add_text_box(slide, "What this means",
                 Inches(0.4), Inches(2.95), Inches(6.0), Inches(0.4),
                 font_size=Pt(15), bold=True, color=C_DARK)
    bullet_box(slide,
        ["We did not fail to detect viruses — we succeeded in distinguishing signal from noise",
         "The pipeline correctly returned zero: a scientifically valid null result at bulk RNA-seq sensitivity",
         "Three \"signals\" from viral-only classification were fully explained by endogenous transcription, k-mer cross-mapping, and index hopping",
         "This establishes a validated baseline: human DRG shows no active exogenous viral transcription above ~10 RPM in these donors",
         "Future cohorts with neuropathic pain patients are now interpretable against this baseline",
        ],
        Inches(0.4), Inches(3.42), Inches(6.0), Inches(3.7),
        font_size=Pt(12.5), gap=6)

    # Right: framing callout box
    rect(slide, Inches(6.7), Inches(2.95), Inches(6.25), Inches(4.2), fill=C_DARK)
    add_text_box(slide, "The sensitivity vs. specificity trade-off — resolved",
                 Inches(6.85), Inches(3.0), Inches(5.9), Inches(0.45),
                 font_size=Pt(14), bold=True, color=C_AMBER)
    add_text_box(slide,
        "Viral-only k-mer classifiers maximize sensitivity by forcing every read into a viral bin. This is correct behavior for environmental metagenomics — where host contamination is the exception.\n\nIn human tissue RNA-seq, the assumption inverts: the vast majority of \"viral\" k-mer matches are host-derived. A classifier with no host reference cannot distinguish endogenous retroviral transcription from exogenous infection.\n\nCompetitive dual-database classification restores specificity without sacrificing sensitivity to true positives. Tier 1 detections — reads that survive competition with the full genome — can be biologically interpreted.",
        Inches(6.85), Inches(3.55), Inches(5.9), Inches(3.45),
        font_size=Pt(12), color=RGBColor(0xEE,0xEE,0xEE), word_wrap=True)


    # ── SLIDE 10: TIER SUMMARY + HEATMAP ────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Dual-database tier classification across all 15 samples",
                 subtitle="Panel A: Tier 2 vs Tier 3 per sample  ·  Panel B: The three Tier 2 taxa — all confirmed false positives",
                 bar_color=C_PURPLE)
    add_figure(slide, FIG3,
               x=Inches(0.35), y=Inches(1.2), w=Inches(12.6))
    footnote(slide, "Tier 2 = viral-only exclusive  ·  Tier 3 = PlusPF-only (non-viral background)  ·  Tier 1 = 0 in all samples (green annotation)")


    # ── SLIDE 11: THE THREE FALSE POSITIVES ──────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "The three Tier 2 taxa — fully resolved mechanisms",
                 subtitle="Each false positive explains a different class of artifact in human tissue metagenomics",
                 bar_color=C_PURPLE)

    fps = [
        {
            "name": "HERV-K (HML-2)\ntaxon 45617",
            "reads": "39,848 reads cohort total",
            "color": C_PURPLE,
            "mech": "Endogenous retroviral transcription",
            "points": [
                "Present in all 15 samples",
                "5.8× enriched in DRG vs. muscle (p < 0.0001 Mann-Whitney U)",
                "Reproduced across 2 independent DRG cohorts at 6 spinal levels",
                "PlusPF: ALL reads reclassified to Homo sapiens",
                "Biological signal: neural-lineage LTR promoter activity in post-mitotic sensory neurons",
                "NOT evidence of exogenous retroviral infection",
            ]
        },
        {
            "name": "Human CMV (HHV-5) [proxy]\ntaxon 3050337",
            "reads": "1,097 reads across 14 samples",
            "color": C_CORAL,
            "mech": "ICTV reclassification k-mer cross-mapping",
            "points": [
                "ICTV 2023 split HHV-5 into primate betaherpesvirus species",
                "LCA rollup routes reads via baboon CMV child node (2169863)",
                "~10× DRG-enriched — mirrors HERV-K (same k-mer source)",
                "PlusPF: ZERO reads at this taxon in ALL samples",
                "Remediation: taxon_remap.tsv → '[proxy]' label; dual-DB confirms no genuine CMV signal",
                "Critical: do NOT interpret as CMV latency in DRG",
            ]
        },
        {
            "name": "Molluscum contagiosum virus\ntaxon 10279",
            "reads": "189 reads total, sporadic",
            "color": C_AMBER,
            "mech": "Index hopping / library prep contamination",
            "points": [
                "Absent from deepest library (Saad_1, 430K bracken_raw reads)",
                "Sporadic: muscle-predominant (65% of reads) — tissue-independent",
                "Distribution inconsistent with k-mer cross-mapping artifact",
                "Consistent with index hopping on patterned flow cell",
                "PlusPF: not detected — no genuine MCV transcription",
                "Does not require follow-up investigation",
            ]
        },
    ]
    for i, fp in enumerate(fps):
        x = Inches(0.3 + i * 4.35)
        rect(slide, x, Inches(1.3), Inches(4.15), Inches(5.8), fill=C_LGRAY)
        rect(slide, x, Inches(1.3), Inches(4.15), Inches(0.72), fill=fp["color"])
        add_text_box(slide, fp["name"],
                     x + Inches(0.1), Inches(1.3), Inches(3.0), Inches(0.7),
                     font_size=Pt(12), bold=True, color=RGBColor(0xFF,0xFF,0xFF))
        add_text_box(slide, fp["reads"],
                     x + Inches(0.1), Inches(2.08), Inches(3.9), Inches(0.3),
                     font_size=Pt(10.5), color=C_DARK, bold=True, italic=True)
        add_text_box(slide, fp["mech"],
                     x + Inches(0.1), Inches(2.38), Inches(3.9), Inches(0.3),
                     font_size=Pt(10.5), color=fp["color"], bold=True)
        bullet_box(slide, fp["points"],
                   x + Inches(0.1), Inches(2.72), Inches(3.9), Inches(3.3),
                   font_size=Pt(11), gap=5)


    # ── SLIDE 12: HERV-K ────────────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "HERV-K — an unexpected biological signal worth pursuing",
                 subtitle="5.8× enrichment in DRG vs. skeletal muscle  ·  reproduced across 2 independent cohorts  ·  6 spinal levels",
                 bar_color=C_PURPLE)
    add_figure(slide, FIG4,
               x=Inches(0.35), y=Inches(1.25), w=Inches(7.6))

    rect(slide, Inches(8.1), Inches(1.3), Inches(5.0), Inches(5.7), fill=C_LGRAY)
    add_text_box(slide, "Why does this matter?",
                 Inches(8.25), Inches(1.38), Inches(4.7), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_DARK)
    bullet_box(slide,
        ["HERV-K/HML-2 is the most recently active HERV family in the human genome",
         "LTR retrotransposons de-repress via NF-κB in neuroinflammatory contexts",
         "ALS literature: TDP-43 pathology correlates with HERV-K envelope protein expression (Li et al. 2015, Science)",
         "MS literature: HERV-K and HERV-W elevated in MS plaques and CSF",
         "DRG-specific enrichment (5.8×) is consistent with neural-lineage LTR transcription in post-mitotic sensory neurons",
         "Hypothesis for Paper 2: does HERV-K expression differ between neuropathic pain vs. control DRG?",
        ],
        Inches(8.25), Inches(1.85), Inches(4.7), Inches(3.8),
        font_size=Pt(12), gap=6)

    rect(slide, Inches(8.1), Inches(6.18), Inches(5.0), Inches(0.75), fill=C_PURPLE)
    add_text_box(slide,
        "Under PlusPF: ALL 39,848 HERV-K reads reclassified to Homo sapiens.\nThis is endogenous chromosomal transcription — not infection.",
        Inches(8.2), Inches(6.22), Inches(4.8), Inches(0.65),
        font_size=Pt(11), bold=True, color=RGBColor(0xFF,0xFF,0xFF))

    footnote(slide, "Mann-Whitney U, DRG vs muscle: p < 0.0001  ·  DRG mean 3,670 reads/sample  ·  Muscle mean 629 reads/sample  ·  donor1_L5 = 7,725 reads (largest DRG library)")


    # ── SECTION DIVIDER 3 ───────────────────────────────────────────────────
    section_divider(prs, 3,
        "The Standard",
        "Where virome-pipeline sits in the field — and how it aligns with best practices",
        color=C_DARK)


    # ── SLIDE 13: TOOL COMPARISON ────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Tool comparison — how virome-pipeline positions in the field",
                 subtitle="Most existing tools were designed for environmental metagenomics or clinical pathogen detection — not bulk RNA-seq from human neural tissue",
                 bar_color=C_DARK)

    headers = ["Tool", "DB type", "Host removal", "FP handling", "HPC-ready", "Neural tissue use"]
    col_ws = [Inches(2.2), Inches(1.6), Inches(1.85), Inches(2.7), Inches(1.35), Inches(3.15)]
    x_starts_t = []
    cx = Inches(0.3)
    for cw in col_ws:
        x_starts_t.append(cx)
        cx += cw + Inches(0.05)

    rows_t = [
        ("virome-pipeline", "Dual-DB: viral + PlusPF", "STAR (splice-aware)", "Dual-DB + 24-entry curated list", "✓ SLURM + Apptainer", "✓ Designed for this"),
        ("MysteryMiner",    "Viral-only k-mer",        "Bowtie2 / custom",    "Prevalence filter only",          "Partial",             "No — general purpose"),
        ("VirCapSeq-VERT",  "Capture + alignment",     "Targeted enrichment", "Alignment-based (high specificity)", "Partial",          "✓ High sensitivity"),
        ("CZ ID (IDseq)",   "NCBI NT + NR",            "STAR + DIAMOND",      "Background correction (ERCC)",    "Cloud only",          "Yes — limited HPC"),
        ("GATK PathSeq",    "Alignment-based",         "Human genome align",  "Coverage + alignment stringency", "✓ GATK framework",    "Yes — designed for it"),
        ("Centrifuge",      "Configurable",            "None (user-adds)",    "None built-in",                   "Yes (fast, low mem)",  "Not neural-specific"),
    ]

    header_y = Inches(1.3)
    for j, (hdr, cw) in enumerate(zip(headers, col_ws)):
        rect(slide, x_starts_t[j], header_y, cw, Inches(0.42), fill=C_DARK)
        add_text_box(slide, hdr,
                     x_starts_t[j] + Inches(0.05), header_y + Inches(0.04),
                     cw - Inches(0.1), Inches(0.35),
                     font_size=Pt(10.5), bold=True, color=RGBColor(0xFF,0xFF,0xFF))

    for i, row in enumerate(rows_t):
        bg = RGBColor(0xE8,0xF4,0xFD) if i == 0 else (C_LGRAY if i % 2 == 1 else RGBColor(0xFF,0xFF,0xFF))
        row_y = Inches(1.75 + i * 0.8)
        for j, (cell, cw) in enumerate(zip(row, col_ws)):
            rect(slide, x_starts_t[j], row_y, cw, Inches(0.76), fill=bg)
            add_text_box(slide, cell,
                         x_starts_t[j] + Inches(0.07), row_y + Inches(0.05),
                         cw - Inches(0.14), Inches(0.68),
                         font_size=Pt(10), color=C_BLUE if i == 0 else C_DARK,
                         bold=(i == 0), word_wrap=True)

    footnote(slide, "VirCapSeq-VERT (Briese et al. 2015) — gold standard for sensitivity but requires targeted capture; not applicable to existing bulk RNA-seq  ·  CZ ID best suited for clinical NGS  ·  PathSeq: alignment-based, slower but high specificity")


    # ── SLIDE 14: TECHNICAL STANDARDS ────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Technical standards — built to production genomics norms",
                 subtitle="Containerization, reproducibility, and HPC integration following nf-core and Nextflow community best practices",
                 bar_color=C_DARK)

    standards = [
        ("Nextflow DSL2", C_BLUE,
         ["Module-based: each step is an independent, re-testable process",
          "Channel-based DAG: data dependencies explicit and visualizable",
          "Aligns with nf-core framework (nf-core/viralrecon, nf-core/taxprofiler)"]),
        ("Containerization\n(Apptainer)", C_BLUE,
         ["Every tool in its own .sif — no version conflicts, no environment drift",
          "Apptainer (formerly Singularity) — rootless, HPC-standard",
          "Reproducibility: same container on any cluster = same result"]),
        ("SLURM integration", C_DARK,
         ["slurm profile: automatic resource requests per process",
          "workDir on /scratch (not groups filesystem) — no quota exhaustion",
          "Interactive compute node workflow prevents login node memory issues"]),
        ("Versioned assets", C_GREEN,
         ["artifact_taxa.tsv: version-controlled, ICTV-aware",
          "taxon_remap.tsv: misleading labels corrected, taxon_id preserved",
          "Samplesheet + config.yaml: full run reproducibility"]),
        ("MultiQC integration", C_AMBER,
         ["Per-sample QC metrics (FastQC, STAR alignment stats, Bracken counts)",
          "Filter funnel and diversity stats injected as custom content",
          "Single HTML report for cohort QC review"]),
        ("Open source (MIT)", C_PURPLE,
         ["github.com/mwilde49/virome-pipeline",
          "DOI via Zenodo (tag v1.2.0 — pending publication)",
          "All dependencies: open source, well-maintained biocontainers"]),
    ]

    for i, (title, color, pts) in enumerate(standards):
        row, col = divmod(i, 3)
        x = Inches(0.3 + col * 4.35)
        y = Inches(1.3 + row * 2.7)
        rect(slide, x, y, Inches(4.15), Inches(2.55), fill=C_LGRAY)
        rect(slide, x, y, Inches(0.2), Inches(2.55), fill=color)
        add_text_box(slide, title,
                     x + Inches(0.28), y + Inches(0.08), Inches(3.75), Inches(0.45),
                     font_size=Pt(13), bold=True, color=color)
        bullet_box(slide, pts,
                   x + Inches(0.28), y + Inches(0.6), Inches(3.75), Inches(1.85),
                   font_size=Pt(11), gap=4)


    # ── SECTION DIVIDER 4 ───────────────────────────────────────────────────
    section_divider(prs, 4,
        "The Ask",
        "Consolidating the lab's sequencing assets — and where we go next",
        color=C_GREEN)


    # ── SLIDE 15: THE ASSET ──────────────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "The lab's sequencing asset — an extraordinary resource",
                 subtitle="100–200 human donors  ·  12+ tissue types  ·  10+ sequencing modalities — most of it not yet systematically analyzed for virome",
                 bar_color=C_GREEN)

    # Left: tissue types
    rect(slide, Inches(0.3), Inches(1.3), Inches(5.5), Inches(5.7), fill=C_LGRAY)
    rect(slide, Inches(0.3), Inches(1.3), Inches(5.5), Inches(0.5), fill=C_GREEN)
    add_text_box(slide, "Tissue types (human, post-mortem or surgical)",
                 Inches(0.4), Inches(1.33), Inches(5.3), Inches(0.43),
                 font_size=Pt(13), bold=True, color=RGBColor(0xFF,0xFF,0xFF))

    tissues = [
        "Dorsal root ganglia (DRG) — primary focus",
        "Spinal cord",
        "Trigeminal ganglia",
        "Vagus nerve",
        "Sympathetic ganglia",
        "Spinal meninges",
        "Nodose ganglia",
        "Adenoids / tonsil",
        "Blood (PBMC, whole blood)",
        "Bronchoalveolar lavage fluid (BAL)",
        "Vastus lateralis (skeletal muscle)",
        "+ additional tissue types",
    ]
    bullet_box(slide, tissues,
               Inches(0.45), Inches(1.88), Inches(5.2), Inches(4.8),
               font_size=Pt(12), gap=5, icon="▸")

    # Middle: modalities
    rect(slide, Inches(6.05), Inches(1.3), Inches(3.4), Inches(5.7), fill=C_LGRAY)
    rect(slide, Inches(6.05), Inches(1.3), Inches(3.4), Inches(0.5), fill=C_BLUE)
    add_text_box(slide, "Modalities available for DRG",
                 Inches(6.15), Inches(1.33), Inches(3.2), Inches(0.43),
                 font_size=Pt(13), bold=True, color=RGBColor(0xFF,0xFF,0xFF))

    modalities = [
        "Bulk RNA-seq  ← virome-pipeline now",
        "ATAC-seq (chromatin)",
        "Iso-seq (long-read)",
        "miRNA-seq",
        "Patch-seq (single cell + electrophys)",
        "PIP-seq (combinatorial barcoding)",
        "scRNA-seq",
        "snRNA-seq  ← pseudobulk virome future",
        "Visium (spatial)  ← pseudobulk virome",
        "Xenium (in situ)  ← cell-type viral tropism",
    ]
    bullet_box(slide, modalities,
               Inches(6.15), Inches(1.88), Inches(3.2), Inches(4.8),
               font_size=Pt(11.5), gap=4, icon="▸")

    # Right: the opportunity
    rect(slide, Inches(9.65), Inches(1.3), Inches(3.4), Inches(5.7), fill=C_DARK)
    add_text_box(slide, "The opportunity",
                 Inches(9.75), Inches(1.38), Inches(3.2), Inches(0.43),
                 font_size=Pt(13), bold=True, color=C_AMBER)
    add_text_box(slide,
        "virome-pipeline is tissue-agnostic — the artifact exclusion list is swappable per tissue type.\n\nA systematic virome screen across all 100–200 donors and all tissue types would be the first comprehensive neural virome atlas of this scale.\n\nAll that is needed:\n\n① Consolidate FASTQs onto a shared location\n② Run a single Nextflow command\n③ Review the MultiQC report\n\nNo new sequencing required for the first atlas pass.",
        Inches(9.75), Inches(1.88), Inches(3.2), Inches(4.8),
        font_size=Pt(12), color=RGBColor(0xDD,0xDD,0xDD), word_wrap=True)


    # ── SLIDE 16: PSEUDOBULK VIROME ──────────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Beyond bulk: pseudobulk virome from snRNA-seq, spatial, and multiome",
                 subtitle="Unmapped reads from single-cell/spatial modalities contain viral signal — with cell-type resolution",
                 bar_color=C_GREEN)

    add_text_box(slide, "The principle",
                 Inches(0.4), Inches(1.32), Inches(5.8), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_DARK)
    bullet_box(slide,
        ["STARsolo / CellRanger generate BAMs with cell barcodes — even for unmapped reads",
         "Unmapped reads can be extracted per cell barcode → pseudobulk per cell type",
         "Kraken2 on these reads = cell-type resolved virome profiling",
         "Key question: if a virus is present, which cell type harbors it? (neurons? satellite glia? macrophages?)",
        ],
        Inches(0.4), Inches(1.8), Inches(5.8), Inches(2.5),
        font_size=Pt(13), gap=6)

    add_text_box(slide, "Available in this lab for DRG",
                 Inches(0.4), Inches(4.45), Inches(5.8), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_DARK)
    bullet_box(slide,
        ["snRNA-seq — pseudobulk per nucleus type",
         "scRNA-seq — pseudobulk per cell type (if low RNA degradation)",
         "Visium — spatial pseudobulk by tissue region",
         "Xenium — single-molecule, cell-type assignment at subcellular resolution",
        ],
        Inches(0.4), Inches(4.9), Inches(5.8), Inches(2.3),
        font_size=Pt(13), gap=6)

    # Right: technical challenges
    rect(slide, Inches(6.5), Inches(1.3), Inches(6.5), Inches(5.7), fill=C_LGRAY)
    add_text_box(slide, "Technical challenges and solutions",
                 Inches(6.65), Inches(1.38), Inches(6.2), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_DARK)
    challenges = [
        ("Fewer reads per cell barcode → lower sensitivity", "Aggregate by cell type (pseudobulk) to recover read depth"),
        ("Host removal was not designed for 10x chemistry", "STAR is already in the pipeline — same index, same logic"),
        ("False positive rate increases at low depth", "Same dual-DB competitive classification; artifact list still applies"),
        ("Cell-type assignment uncertainty", "Use high-confidence clusters only; validate with canonical markers"),
        ("Spatial data: spot contamination", "Kraken2 confidence threshold ≥0.15 for spatial (stricter than bulk)"),
    ]
    for i, (prob, sol) in enumerate(challenges):
        y = Inches(1.9 + i * 0.95)
        rect(slide, Inches(6.65), y, Inches(6.2), Inches(0.88), fill=RGBColor(0xFF,0xFF,0xFF))
        add_text_box(slide, f"⚠  {prob}",
                     Inches(6.75), y + Inches(0.04), Inches(6.0), Inches(0.35),
                     font_size=Pt(11), bold=True, color=C_CORAL)
        add_text_box(slide, f"→  {sol}",
                     Inches(6.75), y + Inches(0.44), Inches(6.0), Inches(0.35),
                     font_size=Pt(11), color=C_GREEN)


    # ── SLIDE 17: CLINICAL COHORT DESIGN ────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Paper 2: clinical cohort design — defining the study with the lab",
                 subtitle="The null result from healthy DRG is the baseline — now we need powered clinical groups",
                 bar_color=C_GREEN)

    add_text_box(slide, "Conditions worth profiling",
                 Inches(0.4), Inches(1.32), Inches(6.0), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_DARK)

    conditions = [
        ("Postherpetic neuralgia (PHN)",        "VZV reactivation in DRG → chronic pain; VZV DNA detectable in DRG by PCR"),
        ("Diabetic peripheral neuropathy (DPN)", "Viral triggers (CMV, EBV) and metabolic inflammation; HERV-K de-repression?"),
        ("Multiple sclerosis",                   "HERV-W/K elevated in MS; demyelination and nerve damage mechanisms"),
        ("HIV-associated neuropathy",            "Direct neurotropism + inflammatory neuropathy; retroviruses in DRG"),
        ("Opioid-associated neuropathy",         "Chronic opioid exposure alters immune surveillance; HERV regulation?"),
        ("Idiopathic small fiber neuropathy",    "Unknown etiology — viral/autoimmune trigger an open question"),
    ]
    for i, (cond, rationale) in enumerate(conditions):
        y = Inches(1.82 + i * 0.84)
        rect(slide, Inches(0.4), y, Inches(2.8), Inches(0.74), fill=C_PURPLE)
        add_text_box(slide, cond,
                     Inches(0.5), y + Inches(0.12), Inches(2.65), Inches(0.55),
                     font_size=Pt(11), bold=True, color=RGBColor(0xFF,0xFF,0xFF), word_wrap=True)
        rect(slide, Inches(3.25), y, Inches(3.0), Inches(0.74), fill=C_LGRAY)
        add_text_box(slide, rationale,
                     Inches(3.35), y + Inches(0.08), Inches(2.8), Inches(0.62),
                     font_size=Pt(10.5), color=C_DARK, word_wrap=True)

    # Right: minimum study design
    rect(slide, Inches(6.5), Inches(1.3), Inches(6.5), Inches(5.7), fill=C_DARK)
    add_text_box(slide, "Minimum study design requirements",
                 Inches(6.65), Inches(1.38), Inches(6.2), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_AMBER)
    reqs = [
        "≥20 DRG donors per group (statistical power for differential abundance)",
        "≥2 clinical groups (neuropathic vs. healthy, or condition-specific)",
        "Tissue-paired samples: DRG + muscle from same surgical specimen",
        "Clinical metadata: diagnosis, age, PMI, tissue handling protocol",
        "Multiple spinal levels per donor where possible (intra-donor replication)",
        "Competitive PlusPF classification (virome-pipeline already implements this)",
        "Alignment-based validation for any signal above detection floor",
    ]
    bullet_box(slide, reqs,
               Inches(6.65), Inches(1.85), Inches(6.2), Inches(4.8),
               font_size=Pt(12), gap=7, icon="✓")


    # ── SLIDE 18: NEXT STEPS + THE ASK ──────────────────────────────────────
    slide = blank_slide(prs)
    slide_header(slide, "Next steps — and what I need from this lab",
                 subtitle="Immediate priorities, the data consolidation project, and the roadmap to a human neural virome atlas",
                 bar_color=C_CORAL)

    # Left column: immediate
    rect(slide, Inches(0.3), Inches(1.3), Inches(3.9), Inches(5.75), fill=C_LGRAY)
    rect(slide, Inches(0.3), Inches(1.3), Inches(3.9), Inches(0.5), fill=C_CORAL)
    add_text_box(slide, "Immediate (this quarter)",
                 Inches(0.4), Inches(1.33), Inches(3.7), Inches(0.43),
                 font_size=Pt(13), bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    bullet_box(slide,
        ["Jayden REJOIN cohort (473-1–473-17): 17 DRG samples — transfer FASTQs to Juno → run → 32 sample cohort",
         "T. gondii BLAST validation (1,330 reads surviving PlusPF — the only unresolved signal)",
         "Add HERV-K (taxon 45617) to artifact list formally",
         "Obtain correct AIG1390 FASTQs from sequencing provider",
         "SRA submission for the 15-sample cohort (paper 1 data availability)",
        ],
        Inches(0.4), Inches(1.88), Inches(3.7), Inches(4.8),
        font_size=Pt(12), gap=7)

    # Middle column: the ask
    rect(slide, Inches(4.45), Inches(1.3), Inches(4.4), Inches(5.75), fill=C_DARK)
    add_text_box(slide, "The ask — data consolidation",
                 Inches(4.55), Inches(1.38), Inches(4.2), Inches(0.43),
                 font_size=Pt(13), bold=True, color=C_AMBER)
    bullet_box(slide,
        ["Inventory all existing bulk RNA-seq samples: which donors, which tissues, which runs are on Juno or local drives?",
         "Identify samples with existing STAR indices / host-removed BAMs — can skip re-alignment",
         "Agree on a central FASTQ storage location on /groups/tprice/",
         "Define a samplesheet standard so any lab member can trigger a virome run",
         "Share sequencing metadata (tissue type, disease status, cohort name) for the combined atlas",
         "Flag samples with paired tissue designs — DRG + muscle from same donor is especially valuable",
        ],
        Inches(4.55), Inches(1.88), Inches(4.2), Inches(4.8),
        font_size=Pt(12), color=RGBColor(0xEE,0xEE,0xEE), gap=7)

    # Right column: roadmap
    rect(slide, Inches(9.1), Inches(1.3), Inches(3.9), Inches(5.75), fill=C_LGRAY)
    rect(slide, Inches(9.1), Inches(1.3), Inches(3.9), Inches(0.5), fill=C_GREEN)
    add_text_box(slide, "Roadmap (6–18 months)",
                 Inches(9.2), Inches(1.33), Inches(3.7), Inches(0.43),
                 font_size=Pt(13), bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    bullet_box(slide,
        ["Paper 1 submission (Bioinformatics Application Note) — pending Zenodo DOI + SRA",
         "32-sample cohort (Jayden) → expand to 50+ with lab bulk RNA-seq inventory",
         "Clinical cohort definition with lab input on pain conditions to prioritize",
         "Pseudobulk virome: snRNA-seq + Visium unmapped reads pipeline",
         "Tissue expansion: spinal cord, TG, vagus nerve, sympathetic ganglia",
         "Paper 2: DRG virome with clinical stratification (Journal of Virology / PLOS Pathogens)",
        ],
        Inches(9.2), Inches(1.88), Inches(3.7), Inches(4.8),
        font_size=Pt(12), gap=7)


    # ── SLIDE 19: CLOSING ────────────────────────────────────────────────────
    slide = blank_slide(prs)
    rect(slide, 0, 0, W, H, fill=C_DARK)
    rect(slide, 0, 0, W, Inches(0.1), fill=C_AMBER)
    rect(slide, 0, H - Inches(0.1), W, Inches(0.1), fill=C_AMBER)

    add_text_box(slide, "Summary",
                 Inches(0.6), Inches(0.5), Inches(4.0), Inches(0.5),
                 font_size=Pt(28), bold=True, color=C_AMBER)

    summary_pts = [
        "virome-pipeline solves the closed-world assumption problem for human neural tissue RNA-seq via competitive dual-database classification",
        "Applied to 15 DRG + muscle samples: zero confirmed viral detections — 100% false positive rate for viral-only Kraken2",
        "Three false positives fully resolved: HERV-K endogenous transcription (5.8× DRG enrichment), CMV k-mer artifact, sporadic contamination",
        "HERV-K tissue differential is a real biological signal worth pursuing — testable hypothesis for neuropathic pain vs. control cohorts",
        "Detection floor of ~10 RPM established — provides power calculation baseline for Paper 2",
        "Pipeline is production-ready, containerized, HPC-deployed, and aligned with nf-core standards",
    ]
    bullet_box(slide, summary_pts,
               Inches(0.6), Inches(1.1), Inches(12.0), Inches(3.8),
               font_size=Pt(13.5), color=RGBColor(0xEE,0xEE,0xEE), gap=8)

    add_text_box(slide, "The vision",
                 Inches(0.6), Inches(5.0), Inches(12.0), Inches(0.45),
                 font_size=Pt(18), bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    add_text_box(slide,
        "A systematic human neural virome atlas — across 12 tissue types, 10 sequencing modalities, and multiple pain/somatosensory conditions — built from the sequencing data this lab already has.",
        Inches(0.6), Inches(5.52), Inches(12.0), Inches(0.9),
        font_size=Pt(15), color=C_AMBER, word_wrap=True)

    add_text_box(slide, "github.com/mwilde49/virome-pipeline  ·  virome-pipeline v1.3.0  ·  MIT License",
                 Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.4),
                 font_size=Pt(11), color=RGBColor(0x88,0x88,0x88))
    add_text_box(slide, "Questions?",
                 Inches(9.5), Inches(5.0), Inches(3.3), Inches(0.8),
                 font_size=Pt(36), bold=True, color=C_GREEN, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    prs = new_prs()
    build(prs)
    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")
